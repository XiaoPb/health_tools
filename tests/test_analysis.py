import json
from collections import Counter
from dataclasses import replace
from pathlib import Path

import numpy as np
import pandas as pd
import pytest

import health_tools.core.analysis.raw as raw_analysis
import health_tools.core.analysis.reporting as analysis_reporting
from health_tools.api import AnalyzeRequest, RequestValidationError, run_analyze
from health_tools.api.analysis_operation import (
    _apply_check_results,
    _apply_compact_check_results,
    _apply_evaluate_results,
    _escalate,
    _generate_psd_plots,
    _generate_raw_plots,
    _offline_records,
    _raw_files,
    _run_supporting_stages,
)
from health_tools.api.context import ExecutionContext
from health_tools.api.models import BatchResult, CheckResult, ItemResult, ItemStatus, OfflineResult
from health_tools.core.analysis.conditions import matches
from health_tools.core.analysis.diagnosis import diagnose
from health_tools.core.analysis.models import AnalysisRecord
from health_tools.core.analysis.psd import analyze_psd_directory
from health_tools.core.analysis.raw import analyze_raw_file, infer_activity
from health_tools.core.analysis.reference import analyze_reference
from health_tools.core.analysis.reporting import (
    _accuracy_rows,
    write_markdown,
    write_ppt,
    write_structured,
)
from health_tools.core.analysis.workspace import AnalysisWorkspace
from health_tools.models.rules import AnalysisRule, ChipRule
from health_tools.rules.loader import RuleLoader

CUSTOM_RULE = """version: '1.0'
type: other
columns:
  reference: REF
  prediction: PRED
  timestamp: time
  ppg_patterns: ['^PPG$']
  acc: [ACCX, ACCY, ACCZ]
detectors: [integrity, raw_signal, reference, accuracy, motion, hr_psd]
sampling: {sample_rate: 10}
thresholds:
  missing_ratio: 0.01
  flat_ratio: 0.98
  saturation_ratio: 0.01
  baseline_drift_ratio: 2
  motion_rms: 0.1
  error: 5
  ref_min: 30
  ref_max: 220
  jump_per_second: 20
causes:
  - id: reference_invalid
    title: reference invalid
    origin: reference
    priority: 100
    when: {feature: reference_valid, op: eq, value: false}
  - id: algorithm_error
    title: algorithm output abnormal
    origin: algorithm
    priority: 50
    when: {feature: algorithm_abnormal, op: eq, value: true}
"""


def _write_csv(path: Path, error: float = 0.0) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    count = 100
    time = np.arange(count) / 10
    frame = pd.DataFrame(
        {
            "time": time,
            "FRAME_ID": np.arange(count),
            "ACCX": np.sin(time) * 0.01,
            "ACCY": np.zeros(count),
            "ACCZ": np.ones(count),
            "PPG": 1000 + np.sin(2 * np.pi * time) * 50,
            "REF": np.full(count, 80.0),
            "PRED": np.full(count, 80.0 + error),
        }
    )
    frame.to_csv(path, index=False)


def test_raw_files_excludes_analysis_auxiliary_csvs(tmp_path: Path):
    source = tmp_path / "data"
    source.mkdir()
    sample = source / "sample.csv"
    sample.write_text("Ipd0,Ipd1\n1,2\n", encoding="utf-8")
    for name in (
        "check_report.csv",
        "check_report_compact.csv",
        "analysis_summary.csv",
        "analysis_diagnosis.csv",
    ):
        (source / name).write_text("file,result\nsample.csv,FAIL\n", encoding="utf-8")

    root, files = _raw_files(source)

    assert root == source
    assert files == [sample]


def test_run_analyze_ignores_nested_output_directory_and_resumes(tmp_path: Path, monkeypatch):
    source = tmp_path / "data"
    source.mkdir()
    sample = source / "sample.csv"
    _write_csv(sample)
    output = source / "analysis_out"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    calls = Counter()

    def fake_run_check_stage(request, source_path, chip, stages, context):
        report = stages / "check" / "check_report.csv"
        report.parent.mkdir(parents=True, exist_ok=True)
        report.write_text("file,result\nsample.csv,PASS\n", encoding="utf-8")
        return report

    def fake_run_raw_stage(request, source_path, rule_object, context, root=None, files=None):
        calls["raw"] += 1
        return (
            [
                AnalysisRecord(
                    file="sample.csv",
                    source=str(sample),
                    analysis_type=request.analysis_type,
                    scene="static",
                    conclusion="未发现异常",
                    confidence=1.0,
                    notes=[f"raw-{calls['raw']}"],
                )
            ],
            source,
            [sample],
            set(),
            request.chip_name,
        )

    monkeypatch.setattr("health_tools.api.analysis_operation.run_check_stage", fake_run_check_stage)
    monkeypatch.setattr("health_tools.api.analysis_operation.run_raw_stage", fake_run_raw_stage)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._apply_check_results", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._apply_compact_check_results",
        lambda *_, **__: None,
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._apply_evaluate_results",
        lambda *_, **__: None,
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_, **__: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_, **__: [])
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots", lambda *_, **__: None
    )

    request = AnalyzeRequest(
        source,
        output,
        analysis_type="other",
        rule_file=str(rule),
        chip_name="gh3036",
        report="markdown",
        allow_offline=False,
    )

    first = run_analyze(request)
    assert _raw_files(source, output)[1] == [sample]

    second = run_analyze(request)

    assert calls["raw"] == 1
    assert first.summary_path == second.summary_path
    assert (output / "stages" / "check" / "check_report.csv").exists()


def test_run_analyze_raw_stage_reuses_discovered_files_when_output_is_nested(
    tmp_path: Path, monkeypatch
):
    source = tmp_path / "data"
    source.mkdir()
    sample = source / "sample.csv"
    _write_csv(sample)
    output = source / "analysis_out"
    stale_stage_csv = output / "stages" / "evaluate" / "file_details.csv"
    stale_stage_csv.parent.mkdir(parents=True)
    stale_stage_csv.write_text("file,mae\nold.csv,99\n", encoding="utf-8")
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")

    monkeypatch.setattr("health_tools.api.analysis_operation.detect_chip", lambda _path: None)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_, **__: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_, **__: [])
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots", lambda *_, **__: None
    )

    result = run_analyze(
        AnalyzeRequest(
            source,
            output,
            analysis_type="other",
            rule_file=str(rule),
            allow_offline=False,
            report="markdown",
        )
    )

    records = json.loads(result.summary_path.read_text(encoding="utf-8"))
    assert [record["file"] for record in records] == ["sample.csv"]
    assert records[0]["source"] == str(sample)


def test_run_analyze_ignores_nested_output_vshb_during_source_detection(
    tmp_path: Path, monkeypatch
):
    source = tmp_path / "data"
    source.mkdir()
    sample = source / "sample.csv"
    _write_csv(sample)
    output = source / "analysis_out"
    stale_vshb = output / "stages" / "offline" / "old_result.vshb"
    stale_vshb.parent.mkdir(parents=True)
    stale_vshb.write_bytes(b"stale")
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")

    monkeypatch.setattr("health_tools.api.analysis_operation.detect_chip", lambda _path: None)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_, **__: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_, **__: [])
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots", lambda *_, **__: None
    )

    result = run_analyze(
        AnalyzeRequest(
            source,
            output,
            analysis_type="other",
            rule_file=str(rule),
            allow_offline=True,
            report="markdown",
        )
    )

    records = json.loads(result.summary_path.read_text(encoding="utf-8"))
    assert [record["file"] for record in records] == ["sample.csv"]
    assert records[0]["source"] == str(sample)


def test_run_analyze_offline_source_excludes_nested_output_vshb(tmp_path: Path, monkeypatch):
    source = tmp_path / "offline"
    source.mkdir()
    real_vshb = source / "real_result.vshb"
    real_vshb.write_text("second,polar,algo_hr,comp_hr,fw_hr\n1,80,80,0,80\n", encoding="utf-8")
    output = source / "analysis_out"
    stale_vshb = output / "stages" / "offline" / "stale_result.vshb"
    stale_vshb.parent.mkdir(parents=True)
    stale_vshb.write_text("second,polar,algo_hr,comp_hr,fw_hr\n1,90,90,0,90\n", encoding="utf-8")
    plotted_vshb = []

    def fake_run_plot(request, *, context=None):
        plotted_vshb.extend(path.name for path in request.input_path.rglob("*_result.vshb"))
        return BatchResult("plot")

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)

    result = run_analyze(AnalyzeRequest(source, output, report="markdown", allow_offline=False))

    records = json.loads(result.summary_path.read_text(encoding="utf-8"))
    assert [record["file"] for record in records] == ["real.csv"]
    assert records[0]["source"] == str(real_vshb)
    assert plotted_vshb == ["real_result.vshb"]


def test_run_analyze_offline_source_with_parent_output_keeps_source_vshb(
    tmp_path: Path, monkeypatch
):
    source = tmp_path / "offline"
    source.mkdir()
    real_vshb = source / "real_result.vshb"
    real_vshb.write_text("second,polar,algo_hr,comp_hr,fw_hr\n1,80,80,0,80\n", encoding="utf-8")
    plotted_vshb = []

    def fake_run_plot(request, *, context=None):
        plotted_vshb.extend(path.name for path in request.input_path.rglob("*_result.vshb"))
        return BatchResult("plot")

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)

    result = run_analyze(AnalyzeRequest(source, tmp_path, report="markdown", allow_offline=False))

    records = json.loads(result.summary_path.read_text(encoding="utf-8"))
    assert [record["file"] for record in records] == ["real.csv"]
    assert records[0]["source"] == str(real_vshb)
    assert plotted_vshb == ["real_result.vshb"]


def test_run_analyze_check_stage_uses_only_discovered_raw_files(tmp_path: Path, monkeypatch):
    source = tmp_path / "data"
    source.mkdir()
    sample = source / "sample.csv"
    _write_csv(sample)
    output = source / "analysis_out"
    for stale in (
        output / "stages" / "evaluate_input" / "evaluate.csv",
        output / "stages" / "offline" / "offline.csv",
    ):
        stale.parent.mkdir(parents=True, exist_ok=True)
        _write_csv(stale)
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    checked_files = []

    def fake_run_check(request, *, context=None):
        assert request.input_path is not None
        checked_files.extend(
            path.relative_to(request.input_path).as_posix()
            for path in sorted(request.input_path.rglob("*.csv"))
        )
        assert request.output_path is not None
        request.output_path.parent.mkdir(parents=True, exist_ok=True)
        request.output_path.write_text(
            "文件相对路径,总异常(结果)\nsample.csv,PASS\n", encoding="utf-8-sig"
        )
        return CheckResult(BatchResult("check"), report_path=request.output_path)

    monkeypatch.setattr("health_tools.api.check_operation.run_check", fake_run_check)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_, **__: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_, **__: [])
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots", lambda *_, **__: None
    )

    run_analyze(
        AnalyzeRequest(
            source,
            output,
            analysis_type="other",
            rule_file=str(rule),
            chip_name="gh3036",
            allow_offline=False,
            report="markdown",
        )
    )

    assert checked_files == ["sample.csv"]


def test_run_analyze_does_not_swallow_check_stage_internal_type_error(tmp_path: Path, monkeypatch):
    source = tmp_path / "input.csv"
    _write_csv(source)
    output = tmp_path / "out"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    calls = Counter()

    def failing_check_stage(request, source_path, chip, stages, context, root, files):
        calls["check"] += 1
        assert root == source.parent
        assert list(files) == [source]
        raise TypeError("check 阶段内部错误")

    monkeypatch.setattr("health_tools.api.analysis_operation.run_check_stage", failing_check_stage)

    with pytest.raises(TypeError, match="check 阶段内部错误"):
        run_analyze(
            AnalyzeRequest(
                source,
                output,
                analysis_type="other",
                rule_file=str(rule),
                chip_name="gh3036",
                allow_offline=False,
                report="markdown",
            )
        )

    assert calls["check"] == 1


def test_run_analyze_supports_legacy_four_argument_raw_stage_override(tmp_path: Path, monkeypatch):
    source = tmp_path / "input.csv"
    _write_csv(source)
    output = tmp_path / "out"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    calls = Counter()

    def legacy_run_raw_stage(request, source_path, rule_object, context):
        calls["raw"] += 1
        return (
            [
                AnalysisRecord(
                    file=source.name,
                    source=str(source),
                    analysis_type=request.analysis_type,
                    scene="static",
                    conclusion="未发现异常",
                    confidence=1.0,
                )
            ],
            source.parent,
            [source],
            set(),
            request.chip_name,
        )

    monkeypatch.setattr("health_tools.api.analysis_operation.run_raw_stage", legacy_run_raw_stage)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_, **__: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_, **__: [])
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots", lambda *_, **__: None
    )
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots", lambda *_, **__: None
    )

    result = run_analyze(
        AnalyzeRequest(
            source,
            output,
            analysis_type="other",
            rule_file=str(rule),
            allow_offline=False,
            report="markdown",
        )
    )

    assert calls["raw"] == 1
    assert result.summary_path.exists()


def test_run_analyze_rejects_directory_with_only_auxiliary_csvs(tmp_path: Path):
    source = tmp_path / "data"
    source.mkdir()
    (source / "check_report.csv").write_text(
        "文件名,总异常(结果)\nsample.csv,FAIL\n", encoding="utf-8"
    )

    with pytest.raises(RequestValidationError, match="未找到可分析的原始 CSV"):
        run_analyze(AnalyzeRequest(source, tmp_path / "out", chip_name="gh3036"))


def test_run_analyze_keeps_missing_rows_from_check_report(tmp_path: Path):
    source = tmp_path / "data"
    existing = source / "a" / "one.csv"
    _write_csv(existing)
    report = source / "check_report.csv"
    pd.DataFrame(
        {
            "文件相对路径": ["a/one.csv", "b/missing.csv"],
            "总异常(结果)": ["PASS", "FAIL"],
        }
    ).to_csv(report, index=False, encoding="utf-8-sig")

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            chip_name="gh3036",
            check_report_path=report,
            allow_offline=False,
        )
    )

    records = json.loads(result.summary_path.read_text(encoding="utf-8"))

    assert [record["file"] for record in records] == ["a/one.csv", "b/missing.csv"]
    missing = records[1]
    assert missing["features"]["input_status"] == "SKIP"
    assert missing["features"]["skip_reason"] == "文件不存在"
    assert missing["notes"] == ["文件不存在"]


def test_run_analyze_rejects_state_request_mismatch_without_restart(tmp_path: Path):
    source = tmp_path / "input.csv"
    _write_csv(source)
    output = tmp_path / "out"
    AnalysisWorkspace.create(output, {"input": "old.csv", "analysis_type": "hr"})

    with pytest.raises(RequestValidationError, match="--restart"):
        run_analyze(
            AnalyzeRequest(
                source,
                output,
                analysis_type="other",
                rule_file=str(tmp_path / "missing.yaml"),
                allow_offline=False,
            )
        )


def test_run_analyze_rejects_no_resume_over_existing_state_without_restart(tmp_path: Path):
    source = tmp_path / "input.csv"
    _write_csv(source)
    output = tmp_path / "out"
    request = AnalyzeRequest(source, output, analysis_type="other", rule_file="same.yaml")
    AnalysisWorkspace.create(
        output,
        {
            "input": str(source),
            "analysis_type": "other",
            "rule": "same.yaml",
        },
    )

    with pytest.raises(RequestValidationError, match="--restart"):
        run_analyze(replace(request, resume=False))


@pytest.mark.parametrize("artifact_name", ["file_diagnosis.csv", "segment_diagnosis.csv"])
def test_run_analyze_rejects_diagnosis_outputs_without_state(artifact_name: str, tmp_path: Path):
    source = tmp_path / "input.csv"
    _write_csv(source)
    output = tmp_path / "out"
    output.mkdir()
    (output / artifact_name).write_text("old", encoding="utf-8")
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")

    with pytest.raises(RequestValidationError, match="分析产物"):
        run_analyze(
            AnalyzeRequest(
                source,
                output,
                analysis_type="other",
                rule_file=str(rule),
                allow_offline=False,
                resume=False,
            )
        )


def test_run_analyze_restart_removes_only_owned_output_artifacts(tmp_path: Path, monkeypatch):
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source)
    external_figure = tmp_path / "external" / "sample.png"
    external_figure.parent.mkdir()
    external_figure.write_bytes(b"png")
    output = tmp_path / "out"
    output.mkdir()
    owned_report = output / "analysis_report.md"
    owned_report.write_text("old", encoding="utf-8")
    owned_stage = output / "stages" / "raw" / "records.json"
    owned_stage.parent.mkdir(parents=True)
    owned_stage.write_text("old", encoding="utf-8")
    workspace = AnalysisWorkspace.create(output, {"input": "old"})
    workspace.start("plot", "old")
    workspace.complete("plot", [external_figure])

    def fail_if_external_removed(*_args, **_kwargs):
        assert external_figure.exists()

    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_raw_plots",
        fail_if_external_removed,
    )

    run_analyze(
        AnalyzeRequest(
            source,
            output,
            analysis_type="other",
            rule_file=str(rule),
            allow_offline=False,
            report="markdown",
            restart=True,
            figure_paths=(external_figure.parent,),
        )
    )

    assert external_figure.exists()
    assert owned_report.read_text(encoding="utf-8").startswith("# PPG 数据分析报告")
    if owned_stage.exists():
        assert owned_stage.read_text(encoding="utf-8") != "old"


def test_run_analyze_resume_after_report_failure_reuses_diagnosis_snapshot(
    tmp_path: Path, monkeypatch
):
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source)
    output = tmp_path / "out"
    calls = Counter()

    def fake_run_raw_stage(request, source_path, rule_object, context, root=None, files=None):
        calls["raw"] += 1
        return (
            [
                AnalysisRecord(
                    "input.csv",
                    str(source),
                    "other",
                    scene="static",
                    conclusion="未发现异常",
                    confidence=1.0,
                    notes=["数据正常"],
                )
            ],
            tmp_path,
            [source],
            set(),
            None,
        )

    def fake_supporting(*_args, **_kwargs):
        calls["evaluate"] += 1
        return [], None, None

    def fake_escalate(*_args, **_kwargs):
        calls["offline"] += 1
        return []

    def fake_plot(*_args, **_kwargs):
        calls["plot"] += 1

    first_report = True

    def flaky_markdown(records, target, accuracy_thresholds=None, accuracy_inclusive=False):
        nonlocal first_report
        calls["report"] += 1
        if first_report:
            first_report = False
            raise RuntimeError("模拟报告生成中断")
        return analysis_reporting.write_markdown(
            records, target, accuracy_thresholds, accuracy_inclusive
        )

    monkeypatch.setattr("health_tools.api.analysis_operation.detect_chip", lambda _path: None)
    monkeypatch.setattr("health_tools.api.analysis_operation.run_raw_stage", fake_run_raw_stage)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages", fake_supporting
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", fake_escalate)
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_psd_plots", fake_plot)
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_raw_plots", fake_plot)
    monkeypatch.setattr("health_tools.api.analysis_operation.write_markdown", flaky_markdown)

    request = AnalyzeRequest(
        source,
        output,
        analysis_type="other",
        rule_file=str(rule),
        allow_offline=False,
        report="markdown",
    )
    with pytest.raises(RuntimeError, match="模拟报告生成中断"):
        run_analyze(request)

    first_counts = calls.copy()
    result = run_analyze(request)

    assert result.reports[0].exists()
    assert calls["report"] == first_counts["report"] + 1
    for stage in ("raw", "evaluate", "offline", "plot"):
        assert calls[stage] == first_counts[stage]


def test_run_analyze_records_all_pipeline_stages(tmp_path: Path, monkeypatch):
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source)

    monkeypatch.setattr("health_tools.api.analysis_operation.detect_chip", lambda _path: None)
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_raw_plots", lambda *_: None)
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_psd_plots", lambda *_: None)

    run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            analysis_type="other",
            rule_file=str(rule),
            allow_offline=False,
            report="markdown",
        )
    )

    state = json.loads((tmp_path / "out" / "analysis_state.json").read_text(encoding="utf-8"))
    statuses = {name: stage["status"] for name, stage in state["stages"].items()}
    assert statuses == {
        "discover": "completed",
        "check": "completed",
        "raw": "completed",
        "evaluate": "completed",
        "offline": "completed",
        "plot": "completed",
        "diagnose": "completed",
        "report": "completed",
    }


def test_run_analyze_invalidates_resume_when_input_file_changes(tmp_path: Path, monkeypatch):
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source)
    output = tmp_path / "out"
    calls = Counter()

    def fake_run_raw_stage(request, source_path, rule_object, context, root=None, files=None):
        calls["raw"] += 1
        return (
            [
                AnalysisRecord(
                    "input.csv",
                    str(source),
                    "other",
                    scene="static",
                    conclusion="未发现异常",
                    confidence=1.0,
                    notes=[f"raw-{calls['raw']}"],
                )
            ],
            tmp_path,
            [source],
            set(),
            None,
        )

    monkeypatch.setattr("health_tools.api.analysis_operation.detect_chip", lambda _path: None)
    monkeypatch.setattr("health_tools.api.analysis_operation.run_raw_stage", fake_run_raw_stage)
    monkeypatch.setattr(
        "health_tools.api.analysis_operation._run_supporting_stages",
        lambda *_args, **_kwargs: ([], None, None),
    )
    monkeypatch.setattr("health_tools.api.analysis_operation._escalate", lambda *_args: [])
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_psd_plots", lambda *_: None)
    monkeypatch.setattr("health_tools.api.analysis_operation._generate_raw_plots", lambda *_: None)

    request = AnalyzeRequest(
        source,
        output,
        analysis_type="other",
        rule_file=str(rule),
        allow_offline=False,
        report="markdown",
    )
    run_analyze(request)
    source.write_text(source.read_text(encoding="utf-8") + "\n", encoding="utf-8")

    run_analyze(request)

    assert calls["raw"] == 2
    summary = json.loads((output / "analysis_summary.json").read_text(encoding="utf-8"))
    assert summary[0]["notes"] == ["raw-2"]


def test_structured_conditions_do_not_evaluate_expressions():
    features = {"value": 5, "flag": True}

    assert matches({"all": [{"feature": "value", "op": "ge", "value": 5}]}, features)
    assert not matches({"feature": "__import__('os')", "op": "eq", "value": True}, features)


def test_algorithm_cause_requires_clean_raw_and_valid_reference():
    rule = AnalysisRule(
        type="other",
        causes=[
            {
                "id": "algorithm",
                "title": "algorithm abnormal",
                "origin": "algorithm",
                "priority": 1,
                "when": {"feature": "algorithm_abnormal", "op": "eq", "value": True},
            }
        ],
    )

    blocked = diagnose(
        {"raw_valid": False, "reference_valid": True, "algorithm_abnormal": True}, rule
    )
    confirmed = diagnose(
        {"raw_valid": True, "reference_valid": True, "algorithm_abnormal": True}, rule
    )

    assert blocked["conclusion"] == "证据不足"
    assert confirmed["conclusion"] == "算法性能极限"
    assert confirmed["actions"] == []


def test_run_analyze_generates_structured_and_markdown_reports(tmp_path: Path):
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source)

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            analysis_type="other",
            rule_file=str(rule),
            report="markdown",
            allow_offline=False,
        )
    )

    assert result.summary_path and result.summary_path.exists()
    file_detail = tmp_path / "out" / "file_diagnosis.csv"
    assert file_detail.exists()
    assert {"within_5", "within_10", "within_15"}.issubset(pd.read_csv(file_detail).columns)
    assert (tmp_path / "out" / "segment_diagnosis.csv").exists()
    report = (tmp_path / "out" / "analysis_report.md").read_text(encoding="utf-8")
    assert "未发现异常" in report
    assert "整体准确度对比" in report
    assert "Online vs Polar" in report
    assert "Offline vs Polar" in report
    assert "±5 bpm" in report
    assert "±10 bpm" in report
    assert "±15 bpm" in report
    assert "异常数据归类" in report
    assert "算法优化" not in report


def test_raw_accuracy_uses_shared_boundary_and_dynamic_strict_thresholds(
    monkeypatch, tmp_path: Path
):
    frame = pd.DataFrame(
        {
            "time": np.arange(7, dtype=float),
            "PPG": 1000 + np.arange(7, dtype=float),
            "REF": [0, 80, 80, 80, 80, 80, 0],
            "PRED": [0, 0, 85, 0, 80, np.nan, 0],
        }
    )
    monkeypatch.setattr(raw_analysis, "_read", lambda path, chip: (frame, None))
    rule = AnalysisRule(
        columns={
            "reference": "REF",
            "prediction": "PRED",
            "timestamp": "time",
            "ppg_patterns": ["^PPG$"],
            "acc": [],
        },
        sampling={"sample_rate": 1},
        thresholds={"error": 5, "ref_min": 30, "ref_max": 220},
    )

    strict, _, _ = analyze_raw_file(
        tmp_path / "sample.csv",
        rule,
        "hr",
        accuracy_thresholds=(5.0, 12.5),
    )
    inclusive, _, _ = analyze_raw_file(
        tmp_path / "sample.csv",
        rule,
        "hr",
        accuracy_thresholds=(5.0, 12.5),
        accuracy_inclusive=True,
    )

    assert strict["metrics"]["samples"] == 3
    assert strict["metrics"]["mae"] == 28.33
    assert strict["metrics"]["within_5"] == 33.33
    assert strict["metrics"]["within_12.5"] == 66.67
    assert inclusive["metrics"]["within_5"] == 66.67
    assert strict["metrics"]["error_ratio"] == pytest.approx(1 / 3)
    assert inclusive["metrics"]["error_ratio"] == strict["metrics"]["error_ratio"]
    assert strict["segments"] == inclusive["segments"]
    assert strict["segments"][0]["start_s"] == 3.0
    assert strict["segments"][0]["end_s"] == 3.0


@pytest.mark.parametrize(
    ("inclusive", "expected"),
    [
        (False, (0.0, 33.33, 66.67)),
        (True, (33.33, 66.67, 100.0)),
    ],
)
def test_raw_accuracy_default_threshold_exact_boundaries(
    monkeypatch, tmp_path: Path, inclusive: bool, expected: tuple[float, float, float]
):
    frame = pd.DataFrame(
        {
            "time": np.arange(3, dtype=float),
            "PPG": [1000.0, 1001.0, 1002.0],
            "REF": [100.0, 100.0, 100.0],
            "PRED": [95.0, 90.0, 85.0],
        }
    )
    monkeypatch.setattr(raw_analysis, "_read", lambda path, chip: (frame, None))
    rule = AnalysisRule(
        columns={
            "reference": "REF",
            "prediction": "PRED",
            "timestamp": "time",
            "ppg_patterns": ["^PPG$"],
            "acc": [],
        },
        sampling={"sample_rate": 1},
        thresholds={"error": 20, "ref_min": 30, "ref_max": 220},
    )

    features, _, _ = analyze_raw_file(
        tmp_path / "sample.csv",
        rule,
        "hr",
        accuracy_inclusive=inclusive,
    )

    metrics = features["metrics"]
    assert (metrics["within_5"], metrics["within_10"], metrics["within_15"]) == expected


def test_accuracy_report_uses_dynamic_thresholds_and_per_comparison_samples(tmp_path: Path):
    records = [
        AnalysisRecord(
            file="first.csv",
            source="first.csv",
            analysis_type="hr",
            scene="dynamic",
            metrics={
                "comparisons": {
                    "online": {
                        "samples": 1,
                        "mae": 10.0,
                        "max_error": 10.0,
                        "within_2": 0.0,
                        "within_7.5": 0.0,
                    },
                    "offline": {
                        "samples": 3,
                        "mae": 0.0,
                        "max_error": 0.0,
                        "within_2": 100.0,
                        "within_7.5": 100.0,
                    },
                }
            },
        ),
        AnalysisRecord(
            file="second.csv",
            source="second.csv",
            analysis_type="hr",
            scene="dynamic",
            metrics={
                "comparisons": {
                    "online": {
                        "samples": 3,
                        "mae": 0.0,
                        "max_error": 0.0,
                        "within_2": 100.0,
                        "within_7.5": 100.0,
                    }
                }
            },
        ),
        AnalysisRecord(
            file="fallback.csv",
            source="fallback.csv",
            analysis_type="hr",
            scene="static",
            metrics={
                "comparisons": {
                    "online_vs_offline": {
                        "samples": 2,
                        "mae": 1.0,
                        "max_error": 1.0,
                        "within_2": 100.0,
                        "within_7.5": 100.0,
                    }
                }
            },
        ),
    ]

    rows = _accuracy_rows(records, (2.0, 7.5))
    overall_online = next(
        row for row in rows if row["scene"] == "整体" and row["comparison"] == "Online vs Polar"
    )
    fallback = next(row for row in rows if row["comparison"] == "Online vs Offline")
    assert overall_online["samples"] == 4
    assert overall_online["mae"] == 2.5
    assert overall_online["within_2"] == 75.0
    assert fallback["samples"] == 2

    structured = write_structured(records, tmp_path / "structured", (2.0, 7.5))
    columns = pd.read_csv(structured[1]).columns
    assert "within_2" in columns
    assert "within_7.5" in columns
    assert "within_5" not in columns
    report = write_markdown(records, tmp_path / "report.md", (2.0, 7.5)).read_text(encoding="utf-8")
    assert "±2 bpm" in report
    assert "±7.5 bpm" in report
    assert "75.0%" in report
    assert "Online vs Offline" in report


def test_accuracy_report_preserves_threshold_input_order(tmp_path: Path):
    thresholds = (15.0, 5.0, 10.0)
    records = [
        AnalysisRecord(
            file="sample.csv",
            source="sample.csv",
            analysis_type="hr",
            metrics={
                "samples": 2,
                "mae": 1.0,
                "max_error": 2.0,
                "within_15": 100.0,
                "within_5": 50.0,
                "within_10": 75.0,
            },
        )
    ]

    structured = write_structured(records, tmp_path / "structured", thresholds)
    columns = pd.read_csv(structured[1]).columns.tolist()
    assert columns[-3:] == ["within_15", "within_5", "within_10"]

    report = write_markdown(records, tmp_path / "report.md", thresholds).read_text(encoding="utf-8")
    header = next(line for line in report.splitlines() if line.startswith("| 对比对象"))
    assert header.index("±15 bpm") < header.index("±5 bpm") < header.index("±10 bpm")


def test_markdown_report_accepts_external_figure_path(tmp_path: Path):
    external_figure = tmp_path / "input" / "figures" / "sample.png"
    external_figure.parent.mkdir(parents=True)
    external_figure.write_bytes(b"png")
    output = tmp_path / "report" / "report.md"
    record = AnalysisRecord(
        file="sample.csv",
        source="sample.csv",
        analysis_type="hr",
        figure=str(external_figure),
    )

    report = write_markdown([record], output, (2.0, 7.5)).read_text(encoding="utf-8")

    assert external_figure.resolve().as_uri() in report


def test_accuracy_report_files_exclude_zero_sample_comparisons():
    records = [
        AnalysisRecord(
            file="empty.csv",
            source="empty.csv",
            analysis_type="hr",
            metrics={
                "comparisons": {
                    "online": {
                        "samples": 0,
                        "mae": 0.0,
                        "max_error": 0.0,
                        "within_5": 0.0,
                    }
                }
            },
        ),
        AnalysisRecord(
            file="valid.csv",
            source="valid.csv",
            analysis_type="hr",
            metrics={
                "comparisons": {
                    "online": {
                        "samples": 2,
                        "mae": 1.0,
                        "max_error": 2.0,
                        "within_5": 100.0,
                    }
                }
            },
        ),
    ]

    rows = _accuracy_rows(records)
    overall_online = next(
        row for row in rows if row["scene"] == "整体" and row["comparison"] == "Online vs Polar"
    )

    assert overall_online["files"] == 1
    assert overall_online["samples"] == 2


def test_accuracy_report_falls_back_to_raw_metrics_for_empty_comparisons():
    record = AnalysisRecord(
        file="sample.csv",
        source="sample.csv",
        analysis_type="hr",
        metrics={
            "samples": 2,
            "mae": 1.0,
            "max_error": 2.0,
            "within_5": 100.0,
            "within_10": 100.0,
            "within_15": 100.0,
            "comparisons": {},
        },
    )

    rows = _accuracy_rows([record])
    overall_online = next(
        row for row in rows if row["scene"] == "整体" and row["comparison"] == "Online vs Polar"
    )

    assert overall_online["available"] is True
    assert overall_online["files"] == 1
    assert overall_online["samples"] == 2
    assert overall_online["mae"] == 1.0


def test_supporting_evaluate_receives_analyze_accuracy_options(monkeypatch, tmp_path: Path):
    source = tmp_path / "input.csv"
    source.write_text("value\n1\n", encoding="utf-8")
    requests = []

    def fake_run_evaluate(request, *, context=None):
        requests.append(request)
        return BatchResult("evaluate")

    monkeypatch.setattr("health_tools.api.operations.run_evaluate", fake_run_evaluate)

    _run_supporting_stages(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            accuracy_thresholds=(2.0, 7.5),
            accuracy_inclusive=True,
        ),
        source,
        [source],
        tmp_path,
        None,
        tmp_path / "stages",
        AnalysisRule(),
        ExecutionContext(),
    )

    assert requests[0].accuracy_thresholds == (2.0, 7.5)
    assert requests[0].accuracy_inclusive is True


@pytest.mark.parametrize(
    ("status", "expected_conclusion"),
    [("FAIL", "原始数据问题"), ("WARNING", "未发现异常")],
)
def test_check_report_changes_diagnosis_only_for_failures(
    tmp_path: Path, status: str, expected_conclusion: str
):
    report = tmp_path / "check_report.csv"
    pd.DataFrame(
        [
            {
                "文件名": "sample.csv",
                "文件相对路径": "dynamic/sample.csv",
                "总异常(结果)": "FAIL" if status == "FAIL" else "PASS",
                "帧完整性(结果)": status,
            }
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    record = AnalysisRecord(
        file="dynamic/sample.csv",
        source=str(tmp_path / "sample.csv"),
        analysis_type="hr",
        features={
            "data_complete": True,
            "raw_valid": True,
            "reference_valid": True,
            "algorithm_abnormal": False,
        },
    )

    _apply_check_results([record], report, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert record.conclusion == expected_conclusion
    assert record.features["check_failures"] == (["帧完整性"] if status == "FAIL" else [])


def test_check_report_prefers_relative_path_for_duplicate_filenames(tmp_path: Path):
    report = tmp_path / "check_report.csv"
    pd.DataFrame(
        [
            {
                "文件名": "x.csv",
                "文件相对路径": "a/x.csv",
                "总异常(结果)": "FAIL",
                "帧完整性(结果)": "FAIL",
                "IPD范围(结果)": "PASS",
            },
            {
                "文件名": "x.csv",
                "文件相对路径": "b/x.csv",
                "总异常(结果)": "FAIL",
                "帧完整性(结果)": "PASS",
                "IPD范围(结果)": "FAIL",
            },
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    records = [
        AnalysisRecord(
            "a/x.csv",
            "x.csv",
            "hr",
            features={
                "data_complete": True,
                "raw_valid": True,
                "reference_valid": True,
                "algorithm_abnormal": False,
            },
        ),
        AnalysisRecord(
            "b/x.csv",
            "x.csv",
            "hr",
            features={
                "data_complete": True,
                "raw_valid": True,
                "reference_valid": True,
                "algorithm_abnormal": False,
            },
        ),
    ]

    _apply_check_results(records, report, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert records[0].features["check_failures"] == ["帧完整性"]
    assert records[1].features["check_failures"] == ["IPD范围"]


def test_check_report_skips_ambiguous_basename_fallback(tmp_path: Path):
    report = tmp_path / "check_report.csv"
    pd.DataFrame(
        [
            {
                "文件名": "x.csv",
                "文件相对路径": "",
                "总异常(结果)": "FAIL",
                "帧完整性(结果)": "FAIL",
            }
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    records = [
        AnalysisRecord("a/x.csv", "x.csv", "hr", features={"raw_valid": True}),
        AnalysisRecord("b/x.csv", "x.csv", "hr", features={"raw_valid": True}),
    ]

    _apply_check_results(records, report, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert "check_failures" not in records[0].features
    assert "check_failures" not in records[1].features


def test_acc_check_evidence_includes_description_and_positions(tmp_path: Path):
    report = tmp_path / "check_report.csv"
    pd.DataFrame(
        [
            {
                "文件名": "sample.csv",
                "文件相对路径": "sample.csv",
                "总异常(结果)": "FAIL",
                "ACC异常(结果)": "FAIL",
                "ACC异常(说明)": "检测到ACC异常帧 16/100 (16.0%)",
                "ACC静止XYZ次数": 16,
                "ACC静止XYZ最长帧": 11,
                "ACC静止XYZ前10帧": "40,52,61,70",
            }
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    record = AnalysisRecord(
        file="sample.csv",
        source=str(tmp_path / "sample.csv"),
        analysis_type="hr",
        features={
            "data_complete": True,
            "raw_valid": True,
            "reference_valid": True,
            "algorithm_abnormal": False,
        },
    )

    _apply_check_results([record], report, RuleLoader.load_analysis_rule("analysis_hr.yaml"))
    evaluate = tmp_path / "file_details.csv"
    pd.DataFrame([{"file": "sample.csv", "mae": 3.0, "samples": 100}]).to_csv(evaluate, index=False)
    _apply_evaluate_results([record], evaluate, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert record.cause and record.cause["id"] == "acc_invalid"
    assert record.cause["title"] == "check 检测到 ACC 静止异常"
    assert "检测到ACC异常帧 16/100" in record.notes[0]
    assert "静止最长连续帧=11" in record.notes[0]
    assert "静止异常帧位置(前10)=40，52，61，70" in record.notes[0]


def test_evaluate_results_prefer_relative_path_for_duplicate_filenames(tmp_path: Path):
    evaluate = tmp_path / "file_details.csv"
    pd.DataFrame(
        [
            {"file": "a/x.csv", "mae": 3.0, "samples": 100},
            {"file": "b/x.csv", "mae": 7.0, "samples": 120},
        ]
    ).to_csv(evaluate, index=False)
    records = [
        AnalysisRecord("a/x.csv", "x.csv", "hr"),
        AnalysisRecord("b/x.csv", "x.csv", "hr"),
    ]

    _apply_evaluate_results(records, evaluate, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert records[0].metrics["evaluate_mae"] == 3.0
    assert records[1].metrics["evaluate_mae"] == 7.0


def test_evaluate_results_skip_ambiguous_basename_fallback(tmp_path: Path):
    evaluate = tmp_path / "file_details.csv"
    pd.DataFrame([{"file": "x.csv", "mae": 3.0, "samples": 100}]).to_csv(evaluate, index=False)
    records = [
        AnalysisRecord("a/x.csv", "x.csv", "hr"),
        AnalysisRecord("b/x.csv", "x.csv", "hr"),
    ]

    _apply_evaluate_results(records, evaluate, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert "evaluate_mae" not in records[0].metrics
    assert "evaluate_mae" not in records[1].metrics


def test_spo2_raw_evidence_uses_plot_ac(monkeypatch, tmp_path: Path):
    source = tmp_path / "sample.csv"
    source.write_text("value\n1\n", encoding="utf-8")
    image = tmp_path / "ac.png"
    image.write_bytes(b"png")
    requests = []

    def fake_run_plot(request, *, context=None):
        requests.append(request)
        return BatchResult(
            "plot",
            (ItemResult(ItemStatus.OK, str(source), str(image)),),
            (image,),
        )

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)
    record = AnalysisRecord(
        file="sample.csv",
        source=str(source),
        analysis_type="spo2",
        focused=True,
        features={"ppg_columns": ["Ipd0", "Ipd1"], "sample_rate": 25},
    )

    _generate_raw_plots(
        AnalyzeRequest(source, tmp_path / "out", analysis_type="spo2"),
        [record],
        "gh3036",
        tmp_path / "figures",
        RuleLoader.load_analysis_rule("analysis_spo2.yaml"),
        ExecutionContext(),
    )

    assert requests[0].plot_type == "ac"
    assert requests[0].channels == "Ipd0,Ipd1"
    assert record.figure == str(image)


def test_spo2_rejects_excessive_motion_even_when_accuracy_is_normal(tmp_path: Path):
    source = tmp_path / "spo2.csv"
    count = 100
    time = np.arange(count) / 25
    pd.DataFrame(
        {
            "TimeStamp": time,
            "FRAME_ID": np.arange(count),
            "ACCX": np.sin(2 * np.pi * time * 2) * 2,
            "ACCY": np.zeros(count),
            "ACCZ": np.ones(count),
            "Ipd0": 1000 + np.sin(2 * np.pi * time) * 50,
            "Ipd1": 900 + np.sin(2 * np.pi * time) * 45,
            "REF_RESULT5": np.full(count, 98.0),
            "ALGO_RESULT5": np.full(count, 98.0),
        }
    ).to_csv(source, index=False)
    rule = RuleLoader.load_analysis_rule("analysis_spo2.yaml")

    result, _, _ = analyze_raw_file(source, rule, "spo2")
    decision = diagnose({**result["features"], **result["metrics"]}, rule)

    assert result["features"]["motion_excessive"] is True
    assert decision["cause"]["id"] == "motion_excessive"
    assert decision["conclusion"] == "原始数据问题"
    assert "静止" in decision["cause"]["title"]


@pytest.mark.parametrize(
    ("column", "expected_min", "expected_max"),
    [("CH0", 5.0, 10.0), ("Rawdata0", 5.0, 10.0), ("Ipd0", 0.0, 1.0)],
)
def test_pi_applies_adc_offset_only_to_rawdata_columns(
    monkeypatch, tmp_path: Path, column: str, expected_min: float, expected_max: float
):
    count = 250
    time = np.arange(count) / 25
    frame = pd.DataFrame(
        {
            "time": time,
            column: 1100 + np.sin(2 * np.pi * time) * 10,
            "REF": np.full(count, 80.0),
            "PRED": np.full(count, 80.0),
        }
    )
    chip_rule = ChipRule(
        chip="test",
        csv={},
        columns=list(frame.columns),
        chip_info={"adc_offset": 1000},
    )
    monkeypatch.setattr(raw_analysis, "_read", lambda path, chip: (frame, chip_rule))
    rule = AnalysisRule(
        columns={
            "reference": "REF",
            "prediction": "PRED",
            "timestamp": "time",
            "ppg_patterns": [f"^{column}$"],
            "acc": [],
        },
        sampling={"sample_rate": 25},
        thresholds={"ref_min": 30, "ref_max": 220},
    )

    result, _, _ = analyze_raw_file(tmp_path / "sample.csv", rule, "hr", chip_name="test")

    assert expected_min < result["features"]["pi"] < expected_max
    expected_unit = "pA" if column.startswith("Ipd") else "adc_lsb"
    assert result["features"]["pi_units"][column] == expected_unit


def test_near_full_only_counts_positive_adc_upper_rail(monkeypatch, tmp_path: Path):
    frame = pd.DataFrame({"Rawdata0": np.full(100, -900.0)})
    chip_rule = ChipRule(
        chip="test",
        csv={},
        columns=["Rawdata0"],
        chip_info={"adc_offset": 1000, "adc_full_scale": 2000},
    )
    monkeypatch.setattr(raw_analysis, "_read", lambda path, chip: (frame, chip_rule))
    rule = AnalysisRule(
        columns={"ppg_patterns": ["^Rawdata0$"], "acc": []},
        sampling={"sample_rate": 25},
    )

    result, _, _ = analyze_raw_file(tmp_path / "sample.csv", rule, "hr", chip_name="test")

    assert result["features"]["near_full_ratio"] == 0.0


def test_spo2_low_pi_requires_recollection(tmp_path: Path):
    source = tmp_path / "spo2.csv"
    count = 250
    time = np.arange(count) / 25
    pd.DataFrame(
        {
            "TimeStamp": time,
            "FRAME_ID": np.arange(count),
            "ACCX": np.zeros(count),
            "ACCY": np.zeros(count),
            "ACCZ": np.full(count, 512),
            "Ipd0": 1000 + np.sin(2 * np.pi * time),
            "Ipd1": 900 + np.sin(2 * np.pi * time) * 0.8,
            "REF_RESULT5": np.full(count, 98.0),
            "ALGO_RESULT5": np.full(count, 98.0),
        }
    ).to_csv(source, index=False)
    rule = RuleLoader.load_analysis_rule("analysis_spo2.yaml")

    result, _, _ = analyze_raw_file(source, rule, "spo2")
    decision = diagnose({**result["features"], **result["metrics"]}, rule)

    assert result["features"]["pi_low"] is True
    assert decision["cause"]["id"] == "low_perfusion"
    assert decision["conclusion"] == "原始数据问题"


def test_dynamic_scene_does_not_use_pi(tmp_path: Path):
    source = tmp_path / "sample.csv"
    _write_csv(source)
    rule = AnalysisRule(
        columns={
            "reference": "REF",
            "prediction": "PRED",
            "timestamp": "time",
            "ppg_patterns": ["^PPG$"],
            "acc": ["ACCX", "ACCY", "ACCZ"],
        },
        sampling={"sample_rate": 10},
        thresholds={"pi_low": 0.5},
    )

    static, _, _ = analyze_raw_file(
        source,
        rule,
        "hr",
        ref_column="REF",
        pred_column="PRED",
        timestamp_column="time",
        scene_override="static",
    )
    dynamic, _, _ = analyze_raw_file(
        source,
        rule,
        "hr",
        ref_column="REF",
        pred_column="PRED",
        timestamp_column="time",
        scene_override="dynamic",
    )

    assert static["features"]["pi"] is not None
    assert dynamic["features"]["pi"] is None
    assert dynamic["features"]["pi_low"] is False


def test_spo2_offline_directory_does_not_enable_hr_psd(tmp_path: Path):
    source = tmp_path / "offline"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n1,98,98,0,98\n",
        encoding="utf-8",
    )

    with pytest.raises(RequestValidationError, match="未启用 hr_psd"):
        run_analyze(
            AnalyzeRequest(
                source,
                tmp_path / "out",
                analysis_type="spo2",
                report="markdown",
                allow_offline=False,
            )
        )


def test_psd_evidence_is_generated_by_plot_command(monkeypatch, tmp_path: Path):
    source = tmp_path / "offline"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n1,100,100,0,100\n",
        encoding="utf-8",
    )
    image = tmp_path / "sample.png"
    image.write_bytes(b"png")
    requests = []

    def fake_run_plot(request, *, context=None):
        requests.append(request)
        return BatchResult(
            "plot",
            (ItemResult(ItemStatus.OK, str(source), str(image)),),
            (image,),
        )

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)
    record = AnalysisRecord(
        file="sample.csv",
        source=str(source / "sample_result.vshb"),
        analysis_type="hr",
        psd={"available": True},
        conclusion="证据不足",
    )

    _generate_psd_plots(source, [record], tmp_path / "figures", ExecutionContext())

    assert requests[0].plot_type == "psd"
    assert record.figure == str(image)


def test_evidence_insufficient_offline_result_still_links_psd_plot(monkeypatch, tmp_path: Path):
    source = tmp_path / "offline"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n1,100,100,0,100\n",
        encoding="utf-8",
    )
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    image = tmp_path / "out" / "figures" / "psd" / "sample.png"
    requests = []

    def fake_run_plot(request, *, context=None):
        requests.append(request)
        image.parent.mkdir(parents=True, exist_ok=True)
        image.write_bytes(b"png")
        return BatchResult(
            "plot",
            (ItemResult(ItemStatus.OK, str(source), str(image)),),
            (image,),
        )

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            analysis_type="other",
            rule_file=str(rule),
            report="markdown",
            allow_offline=False,
        )
    )

    report = result.reports[0].read_text(encoding="utf-8")
    assert result.conclusion_counts["证据不足"] == 1
    assert requests[0].plot_type == "psd"
    assert "figures/psd/sample.png" in report


@pytest.mark.parametrize("version", [None, "specified-version"])
def test_offline_escalation_uses_copy_and_requested_version(
    monkeypatch, tmp_path: Path, version: str
):
    root = tmp_path / "input"
    source = root / "dynamic" / "sample.csv"
    source.parent.mkdir(parents=True)
    source.write_text("value\n1\n", encoding="utf-8")
    requests = []

    def fake_run_offline(request, *, context=None):
        requests.append(request)
        request.output_path.mkdir(parents=True, exist_ok=True)
        return OfflineResult(BatchResult("offline"), output_dir=request.output_path)

    monkeypatch.setattr("health_tools.api.offline_operation.run_offline", fake_run_offline)
    record = AnalysisRecord(
        file="dynamic/sample.csv",
        source=str(source),
        analysis_type="hr",
        focused=True,
    )
    request = AnalyzeRequest(
        root,
        tmp_path / "out",
        offline_version=version,
        accuracy_thresholds=(2.0, 7.5),
        accuracy_inclusive=True,
    )

    escalated = _escalate(
        request,
        [record],
        root,
        "gh3036",
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
        tmp_path / "stages",
        ExecutionContext(),
    )

    copied = tmp_path / "stages" / "offline_input" / "dynamic" / "sample.csv"
    assert escalated == [source]
    assert copied.read_text(encoding="utf-8") == source.read_text(encoding="utf-8")
    assert source.exists()
    assert requests[0].input_path == tmp_path / "stages" / "offline_input"
    assert requests[0].ver == version
    assert requests[0].accuracy_thresholds == (2.0, 7.5)
    assert requests[0].accuracy_inclusive is True


def test_no_offline_records_reason_for_evidence_insufficient(tmp_path: Path):
    source = tmp_path / "sample.csv"
    source.write_text("value\n1\n", encoding="utf-8")
    record = AnalysisRecord(
        file="sample.csv",
        source=str(source),
        analysis_type="hr",
        conclusion="证据不足",
        notes=["缺少关键证据"],
    )

    escalated = _escalate(
        AnalyzeRequest(source, tmp_path / "out", allow_offline=False),
        [record],
        tmp_path,
        "gh3036",
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
        tmp_path / "stages",
        ExecutionContext(),
    )

    assert escalated == []
    assert record.notes[-1] == "已通过 --no-offline 禁用离线 PSD 分析"


def test_focus_supports_directory_glob_and_requires_match(tmp_path: Path):
    source = tmp_path / "data"
    _write_csv(source / "dynamic" / "selected.csv")
    _write_csv(source / "static" / "normal.csv")
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            analysis_type="other",
            rule_file=str(rule),
            focus=("dynamic/*.csv",),
            report="markdown",
            allow_offline=False,
        )
    )

    focused = [item for item in result.batch.items if "selected.csv" in item.input]
    assert len(focused) == 1
    with pytest.raises(RequestValidationError, match="未匹配"):
        run_analyze(
            AnalyzeRequest(
                source,
                tmp_path / "other-out",
                analysis_type="other",
                rule_file=str(rule),
                focus=("missing/**/*.csv",),
                report="markdown",
                allow_offline=False,
            )
        )


def _write_locked_psd(result_dir: Path) -> None:
    result_dir.mkdir(parents=True)
    count = 20
    bins = 128
    peak = 50
    ppg = np.ones((count, bins)) * 0.1
    acc = np.ones((count, bins)) * 0.1
    ppg[:, peak] = 10
    acc[:, peak] = 10
    np.savetxt(result_dir / "sample0.prepsd", ppg, delimiter=",")
    np.savetxt(result_dir / "sample.accrmspsd", acc, delimiter=",")
    pd.DataFrame(
        {
            "second": np.arange(count),
            "polar": np.full(count, 100),
            "algo_hr": np.full(count, 130),
            "fw_hr": np.full(count, 130),
        }
    ).to_csv(result_dir / "sample_result.vshb", index=False)


def test_existing_psd_detects_frequency_lock(tmp_path: Path):
    source = tmp_path / "result"
    _write_locked_psd(source)

    result = run_analyze(
        AnalyzeRequest(source, tmp_path / "out", report="markdown", allow_offline=False)
    )

    assert result.conclusion_counts["算法性能极限"] == 1
    report = result.reports[0].read_text(encoding="utf-8")
    assert "Online vs Polar" in report
    assert "Offline vs Polar" in report
    assert "Comp vs Polar" not in report
    details = analyze_psd_directory(
        source, AnalysisRule(thresholds={"psd_clarity": 3, "psd_lock_hz": 0.2})
    )
    assert details["sample.csv"]["psd_locked"] is True


def test_psd_polar_out_of_range_requires_manual_review(tmp_path: Path):
    source = tmp_path / "result"
    _write_locked_psd(source)
    vshb = source / "sample_result.vshb"
    frame = pd.read_csv(vshb)
    frame["polar"] = 300
    frame["algo_hr"] = 300
    frame["fw_hr"] = 300
    frame.to_csv(vshb, index=False)
    rule = RuleLoader.load_analysis_rule("analysis_hr.yaml")

    details = analyze_psd_directory(source, rule)
    decision = diagnose(details["sample.csv"], rule)

    assert details["sample.csv"]["polar_review_required"] is True
    assert any("范围" in issue for issue in details["sample.csv"]["polar_issues"])
    assert decision["cause"] is None
    assert decision["conclusion"] == "证据不足"
    assert "人工复审" in decision["evidence"]


def test_reference_missing_and_stale_samples_reduce_global_valid_ratio():
    values = np.array([100.0] * 130 + [np.nan] * 10 + [101.0] * 60)
    features, mask = analyze_reference(
        values,
        {
            "ref_min": 30,
            "ref_max": 220,
            "ref_valid_ratio": 0.8,
            "ref_stale_seconds": 120,
        },
        sample_rate=1,
    )

    assert features["polar_review_required"] is True
    assert features["reference_valid"] is False
    assert features["reference_valid_ratio"] < 0.8
    assert not mask[:130].any()


def test_local_polar_issue_warns_without_replacing_global_psd_diagnosis(tmp_path: Path):
    source = tmp_path / "result"
    _write_locked_psd(source)
    vshb = source / "sample_result.vshb"
    frame = pd.read_csv(vshb)
    frame.loc[5, "polar"] = 300
    frame.to_csv(vshb, index=False)

    result = run_analyze(
        AnalyzeRequest(source, tmp_path / "out", report="markdown", allow_offline=False)
    )

    summary = pd.read_json(result.summary_path)
    details = pd.read_csv(tmp_path / "out" / "file_diagnosis.csv")
    report = result.reports[0].read_text(encoding="utf-8")
    assert result.conclusion_counts["算法性能极限"] == 1
    assert summary.iloc[0]["metrics"]["comparisons"]["offline"]["samples"] == 20
    assert "Polar 警告" in details.iloc[0]["warnings"]
    assert "Polar 警告" in report
    assert "跟随运动主频" in report
    assert "![证据图]" in report


def test_existing_psd_includes_comp_comparison_when_nonzero(tmp_path: Path):
    source = tmp_path / "result"
    _write_locked_psd(source)
    vshb = source / "sample_result.vshb"
    frame = pd.read_csv(vshb)
    frame["comp_hr"] = 110
    frame.to_csv(vshb, index=False)

    result = run_analyze(
        AnalyzeRequest(source, tmp_path / "out", report="markdown", allow_offline=False)
    )

    report = result.reports[0].read_text(encoding="utf-8")
    assert "Comp vs Polar" in report


def test_offline_records_preserve_actual_accuracy_metric_keys(monkeypatch, tmp_path: Path):
    source = tmp_path / "result"
    source.mkdir()
    psd_result = {
        "available": True,
        "scene": "static",
        "reference_valid": True,
        "samples": 2,
        "mae": 1.0,
        "max_error": 2.0,
        "error_ratio": 0.0,
        "within_2.5": 100.0,
        "comparisons": {
            "online": {
                "samples": 2,
                "mae": 1.0,
                "max_error": 2.0,
                "within_7.5": 100.0,
            }
        },
    }
    monkeypatch.setattr(
        "health_tools.api.analysis_operation.analyze_psd_directory",
        lambda *_args, **_kwargs: {"sample.csv": psd_result},
    )

    records, _ = _offline_records(
        AnalyzeRequest(source, tmp_path / "out"),
        source,
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
    )

    assert records[0].metrics["within_2.5"] == 100.0
    assert records[0].metrics["comparisons"]["online"]["within_7.5"] == 100.0


def test_vshb_accuracy_remains_available_when_psd_files_are_missing(tmp_path: Path):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n" "1,100,104,0,103\n" "2,100,112,0,108\n",
        encoding="utf-8",
    )

    details = analyze_psd_directory(source, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    sample = details["sample.csv"]
    assert sample["available"] is False
    assert sample["comparisons"]["online"]["mae"] == 5.5
    assert sample["comparisons"]["offline"]["mae"] == 8.0
    assert "comp" not in sample["comparisons"]


@pytest.mark.parametrize(("inclusive", "expected"), [(False, 33.33), (True, 66.67)])
def test_vshb_accuracy_uses_global_boundary_and_keeps_middle_zero(
    tmp_path: Path, inclusive: bool, expected: float
):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n"
        "1,0,0,0,0\n"
        "2,80,0,0,0\n"
        "3,80,80,0,85\n"
        "4,80,0,0,0\n"
        "5,80,80,0,80\n"
        "6,80,nan,0,nan\n"
        "7,0,0,0,0\n",
        encoding="utf-8",
    )

    details = analyze_psd_directory(
        source,
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
        accuracy_thresholds=(5.0, 12.5),
        accuracy_inclusive=inclusive,
    )

    comparisons = details["sample.csv"]["comparisons"]
    assert comparisons["offline"]["samples"] == 3
    assert comparisons["offline"]["within_5"] == 66.67
    assert comparisons["online"]["samples"] == 3
    assert comparisons["online"]["within_5"] == expected
    assert comparisons["online"]["within_12.5"] == 66.67
    assert "comp" not in comparisons


@pytest.mark.parametrize(
    ("inclusive", "expected"),
    [
        (False, (0.0, 33.33, 66.67)),
        (True, (33.33, 66.67, 100.0)),
    ],
)
def test_vshb_accuracy_default_threshold_exact_boundaries(
    tmp_path: Path, inclusive: bool, expected: tuple[float, float, float]
):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n"
        "1,100,95,0,95\n"
        "2,100,90,0,90\n"
        "3,100,85,0,85\n",
        encoding="utf-8",
    )

    details = analyze_psd_directory(
        source,
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
        accuracy_inclusive=inclusive,
    )

    metrics = details["sample.csv"]["comparisons"]["offline"]
    assert (metrics["within_5"], metrics["within_10"], metrics["within_15"]) == expected


def test_vshb_accuracy_enables_comp_and_falls_back_without_polar(tmp_path: Path):
    polar_source = tmp_path / "polar"
    polar_source.mkdir()
    (polar_source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n" "1,100,100,105,100\n" "2,100,100,100,100\n",
        encoding="utf-8",
    )
    fallback_source = tmp_path / "fallback"
    fallback_source.mkdir()
    (fallback_source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n"
        "1,0,0,0,0\n"
        "2,0,100,0,105\n"
        "3,0,100,0,0\n"
        "4,0,100,0,100\n"
        "5,0,0,0,0\n",
        encoding="utf-8",
    )
    rule = RuleLoader.load_analysis_rule("analysis_hr.yaml")

    polar = analyze_psd_directory(polar_source, rule)["sample.csv"]["comparisons"]
    fallback = analyze_psd_directory(fallback_source, rule)["sample.csv"]["comparisons"]

    assert set(polar) == {"offline", "online", "comp"}
    assert polar["comp"]["within_5"] == 50.0
    assert set(fallback) == {"online_vs_offline"}
    assert fallback["online_vs_offline"]["samples"] == 3
    assert fallback["online_vs_offline"]["mae"] == 35.0


def test_vshb_accuracy_uses_online_as_primary_when_offline_is_disabled(tmp_path: Path):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n" "1,100,0,0,130\n" "2,100,0,0,130\n",
        encoding="utf-8",
    )

    sample = analyze_psd_directory(source, RuleLoader.load_analysis_rule("analysis_hr.yaml"))[
        "sample.csv"
    ]

    assert set(sample["comparisons"]) == {"online"}
    assert sample["samples"] == 2
    assert sample["mae"] == 30.0
    assert sample["algorithm_abnormal"] is True


def test_analyze_propagates_dynamic_accuracy_to_psd_plot_and_report(monkeypatch, tmp_path: Path):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n" "1,100,100,0,105\n" "2,100,100,0,100\n",
        encoding="utf-8",
    )
    requests = []

    def fake_run_plot(request, *, context=None):
        requests.append(request)
        return BatchResult("plot")

    monkeypatch.setattr("health_tools.api.file_operations.run_plot", fake_run_plot)

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            report="markdown",
            allow_offline=False,
            accuracy_thresholds=(5.0,),
            accuracy_inclusive=True,
        )
    )

    report = result.reports[0].read_text(encoding="utf-8")
    summary = pd.read_json(result.summary_path)
    assert requests[0].accuracy_thresholds == (5.0,)
    assert requests[0].accuracy_inclusive is True
    assert summary.iloc[0]["metrics"]["within_5"] == 100.0
    assert "±5 bpm" in report
    assert "±10 bpm" not in report
    assert "100.0%" in report


def test_analyze_default_accuracy_keeps_psd_plot_call_signature(monkeypatch, tmp_path: Path):
    source = tmp_path / "result"
    source.mkdir()
    (source / "sample_result.vshb").write_text(
        "second,polar,algo_hr,comp_hr,fw_hr\n" "1,100,100,0,100\n",
        encoding="utf-8",
    )
    calls = []

    def fake_generate_psd_plots(source, records, output, context):
        calls.append((source, records, output, context))

    monkeypatch.setattr(
        "health_tools.api.analysis_operation._generate_psd_plots",
        fake_generate_psd_plots,
    )

    run_analyze(AnalyzeRequest(source, tmp_path / "out", report="markdown", allow_offline=False))

    assert len(calls) == 1


def test_analyze_rejects_invalid_accuracy_thresholds(tmp_path: Path):
    source = tmp_path / "input.csv"
    _write_csv(source)

    with pytest.raises(RequestValidationError, match="准确度阈值不能为空"):
        run_analyze(
            AnalyzeRequest(
                source,
                tmp_path / "out",
                accuracy_thresholds=(),
                allow_offline=False,
            )
        )


def test_ppt_report_uses_packaged_template(tmp_path: Path):
    pytest.importorskip("pptx")
    source = tmp_path / "input.csv"
    rule = tmp_path / "analysis" / "custom.yaml"
    rule.parent.mkdir()
    rule.write_text(CUSTOM_RULE, encoding="utf-8")
    _write_csv(source, error=10)

    result = run_analyze(
        AnalyzeRequest(
            source,
            tmp_path / "out",
            analysis_type="other",
            rule_file=str(rule),
            report="pptx",
            allow_offline=False,
        )
    )

    from pptx import Presentation

    deck = Presentation(str(result.reports[0]))
    assert len(deck.slides) >= 4
    text = "\n".join(
        shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    table_text = "\n".join(
        cell.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    )
    assert "PPG 数据分析报告" in text
    assert "整体" in table_text
    assert "Online vs Polar" in table_text
    assert "Offline vs Polar" in table_text
    assert "±5 bpm" in table_text
    assert "±10 bpm" in table_text
    assert "±15 bpm" in table_text
    assert "异常数据归类" in text
    assert "分析完成" in text


def test_ppt_report_uses_dynamic_accuracy_columns(tmp_path: Path):
    pytest.importorskip("pptx")
    output = write_ppt(
        [
            AnalysisRecord(
                file="sample.csv",
                source="sample.csv",
                analysis_type="hr",
                metrics={
                    "samples": 2,
                    "mae": 2.5,
                    "max_error": 5.0,
                    "within_2": 50.0,
                    "within_7.5": 100.0,
                },
            )
        ],
        tmp_path / "report.pptx",
        (2.0, 7.5),
    )

    from pptx import Presentation

    deck = Presentation(str(output))
    table_text = "\n".join(
        cell.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    )
    assert "±2 bpm" in table_text
    assert "±7.5 bpm" in table_text
    assert "±5 bpm" not in table_text
    assert "50.0%" in table_text


def test_ppt_accuracy_tables_paginate_thresholds_and_rows_with_readable_geometry(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    thresholds = tuple(float(value) for value in range(1, 9))
    records = [
        AnalysisRecord(
            file=f"sample_{index}.csv",
            source=f"sample_{index}.csv",
            analysis_type="hr",
            scene=f"scene_{index}",
            conclusion="未发现异常",
            metrics={
                "samples": 2,
                "mae": 1.0,
                "max_error": 2.0,
                **{f"within_{value}": 100.0 for value in range(1, 9)},
            },
        )
        for index in range(10)
    ]

    output = write_ppt(records, tmp_path / "report.pptx", thresholds)
    deck = Presentation(str(output))
    accuracy_slides = [
        slide
        for slide in deck.slides
        if any(getattr(shape, "text", "").startswith("整体准确度对比") for shape in slide.shapes)
    ]

    assert len(accuracy_slides) > 1
    table_text: list[str] = []
    for slide in accuracy_slides:
        shapes = [shape for shape in slide.shapes if shape.has_table]
        assert len(shapes) == 1
        shape = shapes[0]
        table = shape.table
        assert len(table.columns) <= 10
        assert len(table.rows) <= 7
        assert min(column.width for column in table.columns) >= Inches(0.8)
        assert shape.left + shape.width <= deck.slide_width
        assert shape.top + shape.height <= deck.slide_height
        table_text.extend(cell.text for row in table.rows for cell in row.cells)
    assert {f"±{value} bpm" for value in range(1, 9)}.issubset(table_text)
    assert {f"scene_{index}" for index in range(10)}.issubset(table_text)


def test_ppt_summary_table_is_compact_and_centered(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches

    output = write_ppt(
        [
            AnalysisRecord(
                file="sample.csv",
                source="sample.csv",
                analysis_type="hr",
                scene="dynamic",
                metrics={"samples": 100, "mae": 3.2, "max_error": 12, "error_ratio": 0.1},
                conclusion="未发现异常",
            )
        ],
        tmp_path / "report.pptx",
    )

    from pptx import Presentation

    deck = Presentation(str(output))
    tables = [
        shape.table for slide in list(deck.slides)[1:3] for shape in slide.shapes if shape.has_table
    ]
    assert len(tables) == 2
    for table in tables:
        assert all(row.height <= Inches(0.5) for row in table.rows)
        for row in table.rows:
            for cell in row.cells:
                assert cell.vertical_anchor == MSO_ANCHOR.MIDDLE
                assert cell.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER


def test_ppt_body_uses_scene_filename_slots_and_compact_body_font(tmp_path: Path):
    pytest.importorskip("pptx")
    from importlib.resources import files

    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    output = write_ppt(
        [
            AnalysisRecord(
                file="测试文件.csv",
                source="测试文件.csv",
                analysis_type="hr",
                scene="dynamic",
                focused=True,
                conclusion="未发现异常",
            )
        ],
        tmp_path / "report.pptx",
    )

    deck = Presentation(str(output))
    detail_slide = next(
        slide
        for slide in deck.slides
        if any("测试文件.csv" in getattr(shape, "text", "") for shape in slide.shapes)
    )
    template = Presentation(str(files("health_tools") / "templates" / "analysis_report.pptx"))
    template_slide = template.slides[1]
    title = next(
        shape
        for shape in detail_slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
    )
    filename = next(
        shape
        for shape in detail_slide.shapes
        if shape.name in {"文件名副标题", "文本框 8", "文本占位符 4"}
    )
    body = next(
        shape
        for shape in detail_slide.shapes
        if shape.name == "文本占位符 5"
        or (
            shape.is_placeholder
            and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
            and shape.shape_id != filename.shape_id
        )
    )
    template_shapes = {shape.name: shape for shape in template_slide.shapes}
    assert title.text == "dynamic"
    assert filename.text == "测试文件.csv"
    for shape in (title, filename, body):
        expected = template_shapes[shape.name]
        assert (shape.left, shape.top, shape.width, shape.height) == (
            expected.left,
            expected.top,
            expected.width,
            expected.height,
        )
        assert all(
            run.font.size is None
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
        )
    assert body.text_frame.paragraphs[0].line_spacing == 1


def test_ppt_preserves_current_named_template_slots(monkeypatch, tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    template = Path(analysis_reporting.files("health_tools") / "templates" / "analysis_report.pptx")
    source = Presentation(str(template))
    source_slide = source.slides[1]
    by_name = {shape.name: shape for shape in source_slide.shapes}
    expected_names = {
        "标题 3",
        "文本占位符 4",
        "文本占位符 5",
        "内容占位符 6",
        "内容占位符 1",
        "内容占位符 2",
    }
    if not expected_names.issubset(by_name):
        pytest.skip("当前工作树未携带主工作区最新模板；合入主分支后执行命名槽位验收")
    primary_rect = tuple(
        getattr(by_name["内容占位符 1"], key) for key in ("left", "top", "width", "height")
    )
    secondary_rect = tuple(
        getattr(by_name["内容占位符 2"], key) for key in ("left", "top", "width", "height")
    )
    primary = tmp_path / "primary.png"
    secondary = tmp_path / "secondary.png"
    Image.new("RGB", (640, 360), "white").save(primary)
    Image.new("RGB", (640, 180), "gray").save(secondary)
    output = write_ppt(
        [
            AnalysisRecord(
                "run/a.csv",
                "a.csv",
                "hr",
                activity="run",
                focused=True,
                figure=str(primary),
                secondary_figure=str(secondary),
            )
        ],
        tmp_path / "latest.pptx",
    )
    detail = next(
        slide
        for slide in Presentation(str(output)).slides
        if any(getattr(shape, "text", "") == "跑步" for shape in slide.shapes)
    )
    detail_names = {shape.name: shape for shape in detail.shapes}
    for name in ("标题 3", "文本占位符 4", "文本占位符 5", "内容占位符 6"):
        assert not detail_names[name]._element.xpath(".//a:rPr[@sz]")
    pictures = sorted(
        (shape for shape in detail.shapes if shape.shape_type == 13), key=lambda shape: shape.top
    )
    for picture, rect in zip(pictures, (primary_rect, secondary_rect)):
        left, top, width, height = rect
        assert left <= picture.left and top <= picture.top
        assert picture.left + picture.width <= left + width
        assert picture.top + picture.height <= top + height


def test_ppt_includes_psd_page_for_evidence_insufficient(tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    figure = tmp_path / "sample.png"
    Image.new("RGB", (320, 180), "white").save(figure)
    output = write_ppt(
        [
            AnalysisRecord(
                file="sample.csv",
                source="sample_result.vshb",
                analysis_type="hr",
                conclusion="证据不足",
                notes=["PSD 证据不足"],
                figure=str(figure),
            )
        ],
        tmp_path / "report.pptx",
    )

    deck = Presentation(str(output))
    assert any(shape.shape_type == 13 for slide in deck.slides for shape in slide.shapes)
    assert any(
        getattr(shape, "text", "") == "sample.csv"
        for slide in deck.slides
        for shape in slide.shapes
    )


def test_ppt_places_polar_warning_on_separate_review_page(tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    figure = tmp_path / "sample.png"
    Image.new("RGB", (320, 180), "white").save(figure)
    output = write_ppt(
        [
            AnalysisRecord(
                file="sample.csv",
                source="sample_result.vshb",
                analysis_type="hr",
                conclusion="算法性能极限",
                notes=["算法出值可能跟随运动主频"],
                warnings=["Polar 警告：参考值局部跳变，位置(前10)=5，8，13；需人工复审"],
                figure=str(figure),
            )
        ],
        tmp_path / "report.pptx",
    )

    deck = Presentation(str(output))
    slide_text = [
        "\n".join(getattr(shape, "text", "") for shape in slide.shapes) for slide in deck.slides
    ]
    warning_pages = [text for text in slide_text if "Polar 人工复审警告" in text]
    assert len(warning_pages) == 1
    assert "位置(前10)=5，8，13" in warning_pages[0]
    assert "不作为算法或原始数据错误归因" in warning_pages[0]


def test_report_time_plot_is_generated_separately_and_keeps_psd_primary(monkeypatch, tmp_path):
    from health_tools.api.analysis_operation import _generate_report_time_plots
    from health_tools.core.analysis.models import AnalysisRecord, AnalysisSegment

    source = tmp_path / "sample.csv"
    source.write_text("CH0,REF_RESULT0,ALGO_RESULT0\n1,60,60\n", encoding="utf-8")
    primary = tmp_path / "psd.png"
    primary.write_bytes(b"png")
    calls = []

    class FakePlotter:
        fmt = "png"

        def plot_time(self, frame, target, channels=None, **kwargs):
            calls.append((frame, target, channels, kwargs))
            Path(target).write_bytes(b"time")

    monkeypatch.setattr("health_tools.core.plotter.DataPlotter", FakePlotter)
    record = AnalysisRecord(
        file="sample.csv",
        source=str(source),
        analysis_type="hr",
        focused=True,
        figure=str(primary),
        segments=[AnalysisSegment(2.0, 8.0, 150, max_error=12.0)],
        features={"sample_rate": 25, "ppg_columns": ["CH0"]},
    )

    _generate_report_time_plots([record], tmp_path / "figures" / "time", None)

    assert record.figure == str(primary)
    assert record.secondary_figure is not None
    assert Path(record.secondary_figure).parent == tmp_path / "figures" / "time"
    assert calls[0][2] == ["CH0"]
    assert calls[0][3]["time_range"] == (0.0, 10.0)


@pytest.mark.parametrize(
    ("path", "explicit", "expected"),
    [("run/a.csv", "auto", "run"), ("cycle/a.csv", "auto", "cycle"), ("x.csv", "rest", "rest")],
)
def test_activity_uses_explicit_value_then_path_hint(path: str, explicit: str, expected: str):
    assert infer_activity(Path(path), explicit=explicit) == expected


def test_activity_uses_heart_rate_change_after_path_hints():
    assert (
        infer_activity(
            Path("unknown/a.csv"),
            features={"motion_rms": 0.2, "hr_change_range": 35, "hr_direction_changes": 3},
        )
        == "interval"
    )
    assert (
        infer_activity(
            Path("unknown/b.csv"),
            features={"motion_rms": 0.05, "hr_change_range": 25, "hr_start_end_change": -20},
        )
        == "recovery"
    )


def test_hr_wear_causes_require_combined_evidence():
    from health_tools.core.analysis.diagnosis import diagnose

    rule = RuleLoader.load_analysis_rule(
        str(Path(__file__).parents[1] / "src/health_tools/rules/analysis/analysis_hr.yaml")
    )

    base = {
        "raw_valid": True,
        "reference_valid": True,
        "algorithm_abnormal": False,
        "data_complete": True,
    }
    decision = diagnose(
        base | {"scene": "dynamic", "agc_unstable": True, "baseline_drift": True}, rule
    )
    assert decision["cause"]["id"] == "loose_wear"


def test_compact_check_metrics_feed_diagnosis_features(tmp_path: Path):
    report = tmp_path / "check_report_compact.csv"
    pd.DataFrame(
        [
            {
                "文件名": "a.csv",
                "文件相对路径": "dynamic/a.csv",
                "检查项": "数据居中",
                "状态": "FAIL",
                "通道": "Rawdata3",
                "异常占比": 15.0,
                "近0占比": 2.5,
                "近满量程占比": 12.5,
                "AGC变化次数": 8,
                "AGC有效对数": 10,
                "AGC变化占比": 80.0,
            }
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    record = AnalysisRecord("dynamic/a.csv", "a.csv", "hr", scene="dynamic")

    _apply_compact_check_results(
        [record], report, RuleLoader.load_analysis_rule("analysis_hr.yaml")
    )

    assert record.features["near_full_ratio"] == pytest.approx(12.5)
    assert record.features["agc_change_count"] == 8
    assert record.features["agc_change_ratio"] == pytest.approx(0.8)
    assert record.features["check_channel_metrics"]["Rawdata3"]["abnormal_ratio"] == 15.0


def test_compact_check_prefers_relative_path_for_duplicate_filenames(tmp_path: Path):
    report = tmp_path / "check_report_compact.csv"
    pd.DataFrame(
        [
            {
                "文件名": "x.csv",
                "文件相对路径": "a/x.csv",
                "检查项": "数据居中",
                "状态": "FAIL",
                "通道": "Rawdata0",
                "近满量程占比": 10,
            },
            {
                "文件名": "x.csv",
                "文件相对路径": "b/x.csv",
                "检查项": "数据居中",
                "状态": "FAIL",
                "通道": "Rawdata0",
                "近满量程占比": 20,
            },
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    records = [AnalysisRecord("a/x.csv", "x.csv", "hr"), AnalysisRecord("b/x.csv", "x.csv", "hr")]

    _apply_compact_check_results(records, report, RuleLoader.load_analysis_rule("analysis_hr.yaml"))

    assert records[0].features["near_full_ratio"] == 10
    assert records[1].features["near_full_ratio"] == 20


def test_compact_check_basename_fallback_keeps_all_rows_for_one_file(tmp_path: Path):
    report = tmp_path / "check_report_compact.csv"
    pd.DataFrame(
        [
            {
                "文件名": "x.csv",
                "文件相对路径": "",
                "检查项": "数据居中",
                "状态": "FAIL",
                "通道": "Rawdata0",
                "近满量程占比": 10,
            },
            {
                "文件名": "x.csv",
                "文件相对路径": "",
                "检查项": "数据居中",
                "状态": "FAIL",
                "通道": "Rawdata1",
                "近满量程占比": 20,
            },
        ]
    ).to_csv(report, index=False, encoding="utf-8-sig")
    record = AnalysisRecord("x.csv", "x.csv", "hr")

    _apply_compact_check_results(
        [record], report, RuleLoader.load_analysis_rule("analysis_hr.yaml")
    )

    assert set(record.features["check_channel_metrics"]) == {"Rawdata0", "Rawdata1"}


def test_offline_records_honor_explicit_activity(monkeypatch, tmp_path: Path):
    monkeypatch.setattr(
        "health_tools.api.analysis_operation.analyze_psd_directory",
        lambda *_args, **_kwargs: {
            "a.csv": {
                "available": True,
                "scene": "dynamic",
                "raw_valid": True,
                "reference_valid": True,
                "algorithm_abnormal": False,
            }
        },
    )
    records, _ = _offline_records(
        AnalyzeRequest(tmp_path, tmp_path / "out", activity="run"),
        tmp_path,
        RuleLoader.load_analysis_rule("analysis_hr.yaml"),
    )
    assert records[0].activity == "run"


def test_ppt_uses_activity_and_two_template_picture_areas(tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    primary = tmp_path / "primary.png"
    secondary = tmp_path / "secondary.png"
    Image.new("RGB", (640, 360), "white").save(primary)
    Image.new("RGB", (640, 180), "gray").save(secondary)
    output = write_ppt(
        [
            AnalysisRecord(
                "run/a.csv",
                "a.csv",
                "hr",
                activity="run",
                focused=True,
                conclusion="算法性能极限",
                cause={
                    "id": "frequency_lock",
                    "title": "疑似锁频",
                    "origin": "algorithm",
                    "actions": ["不应展示"],
                },
                notes=["运动主频与算法频率一致"],
                figure=str(primary),
                secondary_figure=str(secondary),
            )
        ],
        tmp_path / "report.pptx",
    )
    deck = Presentation(str(output))
    slide = next(s for s in deck.slides if any(getattr(x, "text", "") == "跑步" for x in s.shapes))
    assert sum(1 for shape in slide.shapes if shape.shape_type == 13) == 2
    text = "\n".join(getattr(shape, "text", "") for shape in slide.shapes)
    assert "不应展示" not in text
    assert "建议：" not in text


@pytest.mark.parametrize(
    ("feature", "cause_id"),
    [
        ("rise_lag", "rise_lag"),
        ("recovery_lag", "recovery_lag"),
        ("output_plateau", "output_plateau"),
    ],
)
def test_algorithm_strategy_causes_never_offer_actions(feature: str, cause_id: str):
    rule = RuleLoader.load_analysis_rule(
        str(Path(__file__).parents[1] / "src/health_tools/rules/analysis/analysis_hr.yaml")
    )
    decision = diagnose(
        {"raw_valid": True, "reference_valid": True, "algorithm_abnormal": True, feature: True},
        rule,
    )
    assert decision["cause"]["id"] == cause_id
    assert decision["actions"] == []
