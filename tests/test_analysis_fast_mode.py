from pathlib import Path

import pytest
from click.testing import CliRunner

from health_tools.api import AnalyzeRequest, RequestValidationError
from health_tools.api.analysis_operation import _fast_report_records
from health_tools.cli import main
from health_tools.core.analysis.reporting import write_ppt


def test_fast_report_matches_prefixed_pngs_and_writes_manifest(tmp_path: Path):
    check = tmp_path / "check_report.csv"
    check.write_text("文件名,结果\nscene/sample.csv,FAIL\n", encoding="utf-8-sig")
    figures = tmp_path / "figures"
    figures.mkdir()
    (figures / "0001_scene_sample.csv.png").write_bytes(b"png")
    (figures / "0002_scene_sample.csv_time.png").write_bytes(b"png")

    records, manifest = _fast_report_records(check, (figures,), tmp_path / "out")

    assert len(records) == 1
    assert records[0].file == "scene/sample.csv"
    assert records[0].figure and records[0].secondary_figure
    assert manifest.exists()
    assert "scene/sample.csv" in manifest.read_text(encoding="utf-8")


def test_fast_report_skips_ambiguous_basename(tmp_path: Path):
    check = tmp_path / "check_report.csv"
    check.write_text("file,result\nsample.csv,FAIL\n", encoding="utf-8")
    first = tmp_path / "one"
    second = tmp_path / "two"
    first.mkdir()
    second.mkdir()
    (first / "sample.csv.png").write_bytes(b"png")
    (second / "sample.csv.png").write_bytes(b"png")

    records, manifest = _fast_report_records(check, (first, second), tmp_path / "out")

    assert records[0].figure is None
    assert "歧义" in manifest.read_text(encoding="utf-8")


def test_fast_ppt_contains_only_check_body_and_images(tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    primary = tmp_path / "sample.csv.png"
    secondary = tmp_path / "sample.csv_time.png"
    Image.new("RGB", (320, 180), "white").save(primary)
    Image.new("RGB", (320, 100), "gray").save(secondary)
    from health_tools.core.analysis.models import AnalysisRecord

    output = write_ppt(
        [
            AnalysisRecord(
                file="sample.csv",
                source="sample.csv",
                analysis_type="hr",
                figure=str(primary),
                secondary_figure=str(secondary),
                warnings=["check: FAIL"],
            )
        ],
        tmp_path / "report.pptx",
        fast_mode=True,
    )
    deck = Presentation(str(output))
    text = "\n".join(
        shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "sample.csv" in text
    assert "整体准确度对比" not in text
    detail = next(
        slide
        for slide in deck.slides
        if any(shape.text == "sample.csv" for shape in slide.shapes if hasattr(shape, "text"))
    )
    assert sum(1 for shape in detail.shapes if shape.shape_type == 13) == 2


def test_fast_ppt_keeps_full_page_order_in_standard_mode(tmp_path: Path):
    pytest.importorskip("pptx")
    from PIL import Image
    from pptx import Presentation

    from health_tools.core.analysis.models import AnalysisRecord

    primary = tmp_path / "sample.csv.png"
    secondary = tmp_path / "sample.csv_time.png"
    Image.new("RGB", (320, 180), "white").save(primary)
    Image.new("RGB", (320, 100), "gray").save(secondary)

    record = AnalysisRecord(
        file="sample.csv",
        source="sample.csv",
        analysis_type="hr",
        scene="dynamic",
        scene_label="跑步",
        activity="run",
        focused=True,
        conclusion="算法性能极限",
        cause={"id": "algorithm_error", "title": "疑似算法问题", "origin": "algorithm"},
        notes=["主证据"],
        warnings=["check: FAIL"],
        figure=str(primary),
        secondary_figure=str(secondary),
        metrics={"samples": 1, "mae": 1.0, "max_error": 2.0},
    )

    output = write_ppt([record], tmp_path / "report.pptx")
    deck = Presentation(str(output))
    texts = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in deck.slides
    ]

    positions = {
        label: next(index for index, text in enumerate(texts) if label in text)
        for label in (
            "PPG 数据分析报告",
            "整体准确度对比",
            "异常数据统计",
            "sample.csv",
            "综合结论",
            "分析完成",
        )
    }
    assert positions["PPG 数据分析报告"] < positions["整体准确度对比"] < positions["异常数据统计"]
    assert (
        positions["异常数据统计"]
        < positions["sample.csv"]
        < positions["综合结论"]
        < positions["分析完成"]
    )


def test_fast_report_cli_requires_check_report_and_figure_dir(tmp_path: Path):
    source = tmp_path / "input.csv"
    source.write_text("CH0\n1\n", encoding="utf-8")
    check_report = tmp_path / "check_report.csv"
    check_report.write_text("文件名,结果\nsample.csv,FAIL\n", encoding="utf-8-sig")

    result = CliRunner().invoke(
        main,
        [
            "analyze",
            "-i",
            str(source),
            "-o",
            str(tmp_path / "out"),
            "--fast-report",
            "--check-report",
            str(check_report),
        ],
    )

    assert result.exit_code != 0
    assert "fast-report" in result.output
    assert "--figure-dir" in result.output


def test_fast_request_requires_check_and_figures(tmp_path: Path):
    from health_tools.api.analysis_operation import _validate

    with pytest.raises(RequestValidationError, match="fast-report"):
        _validate(AnalyzeRequest(tmp_path, tmp_path / "out", fast_report=True))
