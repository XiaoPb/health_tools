"""分析结果的 JSON、CSV、Markdown 和 PPT 输出。"""

from __future__ import annotations

import csv
import json
from importlib.resources import files
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

import matplotlib.pyplot as plt
import numpy as np

from health_tools.core.analysis.models import AnalysisRecord

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
plt.rcParams["axes.unicode_minus"] = False


def _accuracy_thresholds(
    thresholds: Optional[Sequence[float]],
) -> Tuple[float, ...]:
    from health_tools.utils.accuracy import (
        DEFAULT_ACCURACY_THRESHOLDS,
        normalize_accuracy_thresholds,
    )

    return normalize_accuracy_thresholds(thresholds) or DEFAULT_ACCURACY_THRESHOLDS


def _accuracy_keys(thresholds: Optional[Sequence[float]]) -> List[str]:
    from health_tools.utils.accuracy import format_accuracy_threshold

    return [
        f"within_{format_accuracy_threshold(value)}" for value in _accuracy_thresholds(thresholds)
    ]


def _accuracy_rows(
    records: List[AnalysisRecord], accuracy_thresholds: Optional[Sequence[float]] = None
) -> List[Dict[str, Any]]:
    groups: Dict[str, List[AnalysisRecord]] = {"整体": records}
    for record in records:
        groups.setdefault(record.scene, []).append(record)
    labels = {
        "online": "Online vs Polar",
        "offline": "Offline vs Polar",
        "comp": "Comp vs Polar",
        "online_vs_offline": "Online vs Offline",
    }
    accuracy_keys = _accuracy_keys(accuracy_thresholds)
    rows: List[Dict[str, Any]] = []
    for name, items in groups.items():
        comparison_items: Dict[str, List[Dict[str, Any]]] = {}
        for item in items:
            comparisons = item.metrics.get("comparisons")
            if isinstance(comparisons, dict) and comparisons:
                for comparison, metrics in comparisons.items():
                    if isinstance(metrics, dict) and int(metrics.get("samples") or 0) > 0:
                        comparison_items.setdefault(str(comparison), []).append(metrics)
            elif int(item.metrics.get("samples") or 0) > 0:
                comparison_items.setdefault("online", []).append(item.metrics)
        names = ["online", "offline"]
        if comparison_items.get("comp"):
            names.append("comp")
        if comparison_items.get("online_vs_offline"):
            names.append("online_vs_offline")
        for comparison in names:
            metrics_list = comparison_items.get(comparison, [])
            samples = sum(int(metrics.get("samples") or 0) for metrics in metrics_list)
            row: Dict[str, Any] = {
                "comparison": labels[comparison],
                "scene": name,
                "files": len(metrics_list),
                "samples": samples,
                "available": samples > 0,
                "mae": None,
                "max_error": None,
                **{key: None for key in accuracy_keys},
            }
            if samples:
                for metric in ("mae", *accuracy_keys):
                    row[metric] = (
                        sum(
                            float(values.get(metric) or 0) * int(values.get("samples") or 0)
                            for values in metrics_list
                        )
                        / samples
                    )
                row["max_error"] = max(
                    float(values.get("max_error") or 0) for values in metrics_list
                )
            rows.append(row)
    return rows


def _cause_counts(records: List[AnalysisRecord]) -> List[Dict[str, Any]]:
    counts: Dict[tuple, int] = {}
    labels = {"raw": "原始数据", "reference": "参考数据", "algorithm": "算法性能边界"}
    for record in records:
        if not record.abnormal:
            continue
        cause = record.cause or {}
        origin = str(cause.get("origin", "unknown"))
        title = str(cause.get("title") or record.conclusion)
        key = (labels.get(origin, "未确定"), title)
        counts[key] = counts.get(key, 0) + 1
    return [
        {"category": category, "cause": cause, "files": count}
        for (category, cause), count in sorted(counts.items())
    ]


def _plain(record: AnalysisRecord) -> Dict[str, Any]:
    return {
        "file": record.file,
        "source": record.source,
        "analysis_type": record.analysis_type,
        "scene": record.scene,
        "focused": record.focused,
        "features": record.features,
        "metrics": record.metrics,
        "segments": [segment.__dict__ for segment in record.segments],
        "psd": record.psd,
        "cause": record.cause,
        "conclusion": record.conclusion,
        "confidence": record.confidence,
        "notes": record.notes,
        "warnings": record.warnings,
        "figure": record.figure,
    }


def write_evidence_figure(record: AnalysisRecord, output: Path):
    output.parent.mkdir(parents=True, exist_ok=True)
    interval = np.asarray(record.plot_data.get("sample_interval", []), dtype=float)
    interval_time = np.asarray(record.plot_data.get("interval_time", []), dtype=float)
    if not len(interval) or not len(interval_time):
        return None
    fig, axis = plt.subplots(1, 1, figsize=(14, 5))
    axis.plot(interval_time, interval, color="#ef5350", linewidth=0.8)
    axis.axhline(np.nanmedian(interval), color="#212121", linestyle="--", label="中位间隔")
    axis.set_title(f"{record.file} | 采样间隔异常")
    axis.set_ylabel("采样间隔 (s)")
    axis.set_xlabel("时间 (s)")
    axis.legend(loc="upper right")
    axis.grid(alpha=0.2)
    fig.tight_layout()
    fig.savefig(output, dpi=150)
    plt.close(fig)
    record.figure = str(output)
    return output


def write_structured(
    records: Iterable[AnalysisRecord],
    output_dir: Path,
    accuracy_thresholds: Optional[Sequence[float]] = None,
) -> List[Path]:
    records = list(records)
    accuracy_keys = _accuracy_keys(accuracy_thresholds)
    output_dir.mkdir(parents=True, exist_ok=True)
    summary = output_dir / "analysis_summary.json"
    summary.write_text(
        json.dumps([_plain(record) for record in records], ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    detail = output_dir / "file_diagnosis.csv"
    with detail.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=[
                "file",
                "scene",
                "focused",
                "conclusion",
                "confidence",
                "cause",
                "evidence",
                "warnings",
                "actions",
                "mae",
                "max_error",
                "error_ratio",
                *accuracy_keys,
            ],
        )
        writer.writeheader()
        for record in records:
            cause = record.cause or {}
            writer.writerow(
                {
                    "file": record.file,
                    "scene": record.scene,
                    "focused": record.focused,
                    "conclusion": record.conclusion,
                    "confidence": round(record.confidence, 3),
                    "cause": cause.get("title", ""),
                    "evidence": record.notes[0] if record.notes else "",
                    "warnings": "；".join(record.warnings),
                    "actions": "；".join(cause.get("actions", [])),
                    "mae": record.metrics.get("mae", ""),
                    "max_error": record.metrics.get("max_error", ""),
                    "error_ratio": record.metrics.get("error_ratio", ""),
                    **{key: record.metrics.get(key, "") for key in accuracy_keys},
                }
            )
    segments = output_dir / "segment_diagnosis.csv"
    with segments.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["file", "start_s", "end_s", "samples", "mean_error", "max_error"],
        )
        writer.writeheader()
        for record in records:
            for segment in record.segments:
                writer.writerow({"file": record.file, **segment.__dict__})
    return [summary, detail, segments]


def write_markdown(
    records: Iterable[AnalysisRecord],
    output: Path,
    accuracy_thresholds: Optional[Sequence[float]] = None,
    accuracy_inclusive: bool = False,
) -> Path:
    from health_tools.utils.accuracy import format_accuracy_threshold

    records = list(records)
    thresholds = _accuracy_thresholds(accuracy_thresholds)
    accuracy_keys = _accuracy_keys(thresholds)
    output.parent.mkdir(parents=True, exist_ok=True)
    lines = ["# PPG 数据分析报告", "", f"文件数：{len(records)}", ""]
    counts: Dict[str, int] = {}
    for record in records:
        counts[record.conclusion] = counts.get(record.conclusion, 0) + 1
    lines.extend(
        ["## 批次结论", "", "；".join(f"{key}: {value}" for key, value in counts.items()), ""]
    )
    lines.extend(
        [
            "## 整体准确度对比",
            "",
            "绝对误差"
            + ("不超过" if accuracy_inclusive else "小于")
            + "对应阈值的有效样本占比如下。",
            "",
            "| 对比对象 | 场景 | 文件数 | 样本数 | MAE | 最大误差 | "
            + " | ".join(f"±{format_accuracy_threshold(value)} bpm" for value in thresholds)
            + " |",
            "|---|---|---:|---:|---:|---:|" + "---:|" * len(thresholds),
        ]
    )
    for row in _accuracy_rows(records, thresholds):
        if row["available"]:
            values = (
                f"{row['mae']:.2f}",
                f"{row['max_error']:.2f}",
                *(f"{row[key]:.1f}%" for key in accuracy_keys),
            )
            files = str(row["files"])
            samples = str(row["samples"])
        else:
            files, samples = "0", "-"
            values = ("未执行", "-", *("-" for _ in accuracy_keys))
        lines.append(
            f"| {row['comparison']} | {row['scene']} | {files} | {samples} | "
            + " | ".join(values)
            + " |"
        )
    lines.extend(
        [
            "",
            "## 异常数据归类",
            "",
            "| 来源 | 具体原因 | 文件数 |",
            "|---|---|---:|",
        ]
    )
    cause_rows = _cause_counts(records)
    if cause_rows:
        for row in cause_rows:
            lines.append(f"| {row['category']} | {row['cause']} | {row['files']} |")
    else:
        lines.append("| - | 未发现异常数据 | 0 |")
    lines.append("")
    for record in records:
        lines.extend(
            [
                f"## {record.file}",
                "",
                f"- 场景：{record.scene}",
                f"- 结论：{record.conclusion}",
                f"- 置信度：{record.confidence:.0%}",
            ]
        )
        if record.cause:
            lines.append(f"- 可能原因：{record.cause.get('title', '')}")
        if record.notes:
            lines.append(f"- 证据：{record.notes[0]}")
        if record.warnings:
            lines.append(f"- 警告：{'；'.join(record.warnings)}")
        actions = (record.cause or {}).get("actions", [])
        if actions:
            lines.append(f"- 原始数据措施：{'；'.join(actions)}")
        if record.figure:
            rel = Path(record.figure).relative_to(output.parent).as_posix()
            lines.extend(["", f"![证据图]({rel})"])
        lines.append("")
    output.write_text("\n".join(lines), encoding="utf-8")
    return output


def _duplicate_slide(prs, source):
    from copy import deepcopy

    slide = _add_slide_before_ending(prs, source.slide_layout)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")
    return slide


def _add_slide_before_ending(prs, layout):
    slide = prs.slides.add_slide(layout)
    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(len(prs.slides._sldIdLst) - 1, slide_id)
    return slide


def _text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _set_text(shape, text: str) -> None:
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text


def _replace_text_preserving_style(shape, text: str) -> None:
    paragraphs = shape.text_frame.paragraphs
    if paragraphs and paragraphs[0].runs:
        paragraphs[0].runs[0].text = text
        for run in paragraphs[0].runs[1:]:
            run.text = ""
        for paragraph in paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
    else:
        _set_text(shape, text)


def _placeholder(slide, placeholder_type):
    return next(
        (
            shape
            for shape in slide.shapes
            if shape.is_placeholder and shape.placeholder_format.type == placeholder_type
        ),
        None,
    )


def _filename_shape(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER

    named = next(
        (
            shape
            for shape in slide.shapes
            if shape.name in {"文件名副标题", "文本框 8", "文本占位符 4"}
        ),
        None,
    )
    if named is not None:
        return named
    subtitle = _placeholder(slide, PP_PLACEHOLDER.SUBTITLE)
    if subtitle:
        return subtitle
    return None


def _body_shape(slide):
    named = next(
        (shape for shape in slide.shapes if shape.name == "文本占位符 5"),
        None,
    )
    if named is not None:
        return named
    from pptx.enum.shapes import PP_PLACEHOLDER

    return _placeholder(slide, PP_PLACEHOLDER.BODY)


def _check_shape(slide):
    return next(
        (shape for shape in slide.shapes if shape.name == "内容占位符 6"),
        None,
    )


def _primary_picture(slide):
    named = next(
        (
            shape
            for shape in slide.shapes
            if shape.name in {"内容占位符 1", "内容占位符 3", "主图占位符"}
        ),
        None,
    )
    if named is not None:
        return named
    from pptx.enum.shapes import PP_PLACEHOLDER

    objects = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder
        and shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT
        and shape.name not in {"内容占位符 6", "内容占位符 2", "副图占位符"}
    ]
    return min(objects, key=lambda shape: shape.top) if objects else None


def _remove_secondary_picture(slide) -> None:
    from pptx.enum.shapes import PP_PLACEHOLDER

    secondary = next(
        (shape for shape in slide.shapes if shape.name == "副图占位符"),
        None,
    )
    if secondary is None:
        secondary = _placeholder(slide, PP_PLACEHOLDER.PICTURE)
    if secondary:
        secondary._element.getparent().remove(secondary._element)


def _secondary_picture(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER

    named = next(
        (shape for shape in slide.shapes if shape.name in {"副图占位符", "内容占位符 2"}),
        None,
    )
    if named is not None:
        return named
    objects = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT
    ]
    return max(objects, key=lambda shape: shape.top) if len(objects) > 1 else None


def _set_compact_body_text(shape, text: str) -> None:
    _set_text(shape, text)
    for paragraph in shape.text_frame.paragraphs:
        paragraph.line_spacing = 1
        paragraph.space_before = 0
        paragraph.space_after = 0


def _add_picture(slide, placeholder, image_path: str) -> None:
    from PIL import Image

    left, top, width, height = (
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    )
    placeholder._element.getparent().remove(placeholder._element)
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    area_ratio = width / height
    if image_ratio >= area_ratio:
        picture = slide.shapes.add_picture(image_path, left, top, width=width)
        picture.top = top + (height - picture.height) // 2
    else:
        picture = slide.shapes.add_picture(image_path, left, top, height=height)
        picture.left = left + (width - picture.width) // 2


def _populate_content(
    slide,
    title_text: str,
    body_text: str,
    figure: str = "",
    subtitle_text: str = "",
    secondary_figure: str = "",
    check_text: str = "",
) -> None:
    from pptx.enum.shapes import PP_PLACEHOLDER

    title = _placeholder(slide, PP_PLACEHOLDER.TITLE)
    body = _body_shape(slide)
    content = _primary_picture(slide)
    subtitle = _filename_shape(slide)
    secondary = _secondary_picture(slide)
    check_shape = _check_shape(slide)
    if title:
        _set_text(title, title_text)
    if subtitle:
        _set_text(subtitle, subtitle_text)
    if body:
        _set_compact_body_text(body, body_text)
    if check_shape:
        _set_text(check_shape, check_text or "未发现 check 异常")
    if content:
        if figure:
            _add_picture(slide, content, figure)
        else:
            content._element.getparent().remove(content._element)
    if secondary:
        if secondary_figure:
            _add_picture(slide, secondary, secondary_figure)
        else:
            secondary._element.getparent().remove(secondary._element)


def _populate_warning(slide, record: AnalysisRecord) -> None:
    from pptx.enum.shapes import PP_PLACEHOLDER

    title = _placeholder(slide, PP_PLACEHOLDER.TITLE)
    body = _body_shape(slide)
    content = _placeholder(slide, PP_PLACEHOLDER.OBJECT)
    subtitle = _filename_shape(slide)
    if title:
        _set_text(title, "Polar 人工复审警告")
    if body:
        _set_compact_body_text(
            body,
            f"文件：{record.file}\n\n"
            "Polar 可能仅在局部异常。\n\n"
            "原分析结论与关键图表保留。\n\n"
            "警告不作为算法或原始数据错误归因。",
        )
    if subtitle:
        _set_text(subtitle, record.file)
    if content:
        _set_text(content, "\n\n".join(record.warnings))
    _remove_secondary_picture(slide)


def _populate_summary(slide, records: List[AnalysisRecord]) -> str:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    causes = _cause_counts(records)
    body_lines = [f"{row['category']}：{row['cause']}（{row['files']}）" for row in causes] or [
        "未发现异常数据"
    ]
    title = _placeholder(slide, PP_PLACEHOLDER.TITLE)
    body = _body_shape(slide)
    content = _placeholder(slide, PP_PLACEHOLDER.OBJECT)
    subtitle = _filename_shape(slide)
    if title:
        _set_text(title, "批次分析结论")
    if body:
        _set_text(body, "异常数据归类\n" + "\n".join(body_lines))
    if subtitle:
        _set_text(subtitle, "")
    if content:
        left, top, width = content.left, content.top, content.width
        content._element.getparent().remove(content._element)
        row_height = Inches(0.42)
        cause_rows = causes or [{"category": "-", "cause": "未发现异常数据", "files": 0}]
        table_height = row_height * (len(cause_rows) + 1)
        table = slide.shapes.add_table(len(cause_rows) + 1, 3, left, top, width, table_height).table
        column_ratios = (0.22, 0.63, 0.15)
        for column, ratio in zip(table.columns, column_ratios):
            column.width = int(width * ratio)
        for row in table.rows:
            row.height = row_height
        headers = ["来源", "具体原因", "文件数"]
        for column, header in enumerate(headers):
            cell = table.cell(0, column)
            cell.text = header
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(13, 110, 253)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(11)
        for row_index, row in enumerate(cause_rows, 1):
            values = [row["category"], row["cause"], str(row["files"])]
            for column, value in enumerate(values):
                cell = table.cell(row_index, column)
                cell.text = value
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.04)
                cell.margin_right = Inches(0.04)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                if row_index % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(238, 246, 255)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"
                        run.font.size = Pt(11)
    _remove_secondary_picture(slide)
    return "\n".join(body_lines)


def _populate_accuracy(
    slide,
    rows: List[Dict[str, Any]],
    thresholds: Sequence[float],
    accuracy_inclusive: bool = False,
    page_number: int = 1,
    page_count: int = 1,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    from health_tools.utils.accuracy import format_accuracy_threshold

    title = _placeholder(slide, PP_PLACEHOLDER.TITLE)
    content = _placeholder(slide, PP_PLACEHOLDER.OBJECT)
    body = _body_shape(slide)
    subtitle = _filename_shape(slide)
    if title:
        title_text = "整体准确度对比"
        if page_count > 1:
            title_text += f"（{page_number}/{page_count}）"
        _set_text(title, title_text)
    thresholds = tuple(thresholds)
    accuracy_keys = _accuracy_keys(thresholds)
    threshold_labels = " / ".join(f"±{format_accuracy_threshold(value)}" for value in thresholds)
    if body:
        relation = "不超过" if accuracy_inclusive else "小于"
        _set_text(body, f"{threshold_labels} bpm 为绝对误差{relation}对应阈值的有效样本占比。")
    if subtitle:
        _set_text(subtitle, "")
    if not content:
        _remove_secondary_picture(slide)
        return
    left, top, width = content.left, content.top, content.width
    content._element.getparent().remove(content._element)
    row_height = Inches(0.46)
    table = slide.shapes.add_table(
        len(rows) + 1,
        6 + len(thresholds),
        left,
        top,
        width,
        row_height * (len(rows) + 1),
    ).table
    fixed_ratios = (0.20, 0.11, 0.08, 0.10, 0.10, 0.11)
    remaining = max(1.0 - sum(fixed_ratios), 0.0)
    threshold_ratio = remaining / len(thresholds)
    column_ratios = (*fixed_ratios, *(threshold_ratio for _ in thresholds))
    for column, ratio in zip(table.columns, column_ratios):
        column.width = int(width * ratio)
    for row in table.rows:
        row.height = row_height
    headers = [
        "对比对象",
        "场景",
        "文件数",
        "样本数",
        "MAE",
        "最大误差",
        *(f"±{format_accuracy_threshold(value)} bpm" for value in thresholds),
    ]
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.03)
        cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(13, 110, 253)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(10)
    for row_index, row in enumerate(rows, 1):
        if row["available"]:
            values = [
                row["comparison"],
                row["scene"],
                str(row["files"]),
                str(row["samples"]),
                f"{row['mae']:.2f}",
                f"{row['max_error']:.2f}",
                *(f"{row[key]:.1f}%" for key in accuracy_keys),
            ]
        else:
            values = [
                row["comparison"],
                row["scene"],
                "0",
                "-",
                "未执行",
                "-",
                *("-" for _ in accuracy_keys),
            ]
        for column, value in enumerate(values):
            cell = table.cell(row_index, column)
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(238, 246, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(10.5)
    _remove_secondary_picture(slide)


def write_ppt(
    records: Iterable[AnalysisRecord],
    output: Path,
    accuracy_thresholds: Optional[Sequence[float]] = None,
    accuracy_inclusive: bool = False,
) -> Path:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError as exc:
        raise RuntimeError("生成 PPT 需要安装 python-pptx") from exc
    template = Path(files("health_tools") / "templates" / "analysis_report.pptx")
    prs = Presentation(str(template))
    records = list(records)
    cover = prs.slides[0]
    content = prs.slides[1]
    ending = prs.slides[-1]
    for shape in _text_shapes(cover):
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            _set_text(shape, "PPG 数据分析报告")
        elif shape.is_placeholder:
            _set_text(shape, "原始数据、准确度与证据归因")
    for shape in _text_shapes(ending):
        if "THANK" in shape.text.upper():
            _replace_text_preserving_style(shape, "分析完成")
    counts: Dict[str, int] = {}
    for record in records:
        counts[record.conclusion] = counts.get(record.conclusion, 0) + 1
    summary_text = "\n".join(f"{name}：{count} 个文件" for name, count in counts.items())
    from pptx.util import Inches

    thresholds = _accuracy_thresholds(accuracy_thresholds)
    rows = _accuracy_rows(records, thresholds)
    threshold_pages = [thresholds[index : index + 3] for index in range(0, len(thresholds), 3)]
    accuracy_slide = _add_slide_before_ending(prs, prs.slide_layouts[5])
    accuracy_content = _placeholder(accuracy_slide, PP_PLACEHOLDER.OBJECT)
    row_height = Inches(0.46)
    rows_per_page = (
        min(max(int(accuracy_content.height // row_height) - 1, 1), 6) if accuracy_content else 1
    )
    row_pages = [
        rows[index : index + rows_per_page] for index in range(0, len(rows), rows_per_page)
    ]
    pages = [
        (row_page, threshold_page) for row_page in row_pages for threshold_page in threshold_pages
    ]
    for page_index, (page_rows, page_thresholds) in enumerate(pages):
        slide = (
            accuracy_slide
            if page_index == 0
            else _add_slide_before_ending(prs, prs.slide_layouts[5])
        )
        _populate_accuracy(
            slide,
            page_rows,
            page_thresholds,
            accuracy_inclusive,
            page_index + 1,
            len(pages),
        )
    detail_records = [
        record
        for record in records
        if record.abnormal
        or record.focused
        or bool(record.warnings)
        or (record.conclusion == "证据不足" and bool(record.figure))
    ]
    activity_names = {
        "rest": "静息",
        "walk": "步行",
        "run": "跑步",
        "cycle": "骑行",
        "strength": "力量训练",
        "interval": "间歇训练",
        "recovery": "恢复阶段",
        "other": "其他场景",
    }
    for record in detail_records:
        slide = _duplicate_slide(prs, content)
        cause = (record.cause or {}).get("title", "无")
        origin = (record.cause or {}).get("origin")
        actions = [] if origin == "algorithm" else list((record.cause or {}).get("actions", []))
        evidence = record.notes[0] if record.notes else "无"
        body_lines = [
            f"结论：{cause if cause != '无' else record.conclusion}",
            f"置信度：{record.confidence:.0%}",
            f"关键证据：{evidence}",
        ]
        if actions:
            body_lines.append(f"建议：{'；'.join(actions)}")
        check_metrics = record.features.get("check_channel_metrics", {})
        check_lines = []
        if isinstance(check_metrics, dict):
            for channel, metrics in check_metrics.items():
                if not isinstance(metrics, dict):
                    continue
                ratio = metrics.get("abnormal_ratio")
                if ratio is not None:
                    check_lines.append(f"{channel} 异常占比 {float(ratio):.1f}%")
        _populate_content(
            slide,
            (
                activity_names.get(record.activity, record.activity)
                if record.activity and record.activity != "other"
                else record.scene
            ),
            "\n".join(body_lines),
            record.figure or "",
            subtitle_text=record.file,
            secondary_figure=record.secondary_figure or "",
            check_text="\n".join(check_lines),
        )
        if record.warnings:
            warning_slide = _duplicate_slide(prs, content)
            _populate_warning(warning_slide, record)
    conclusion_slide = _duplicate_slide(prs, content)
    cause_summary = _populate_summary(content, records)
    _populate_content(
        conclusion_slide,
        "综合结论",
        summary_text + "\n\n" + cause_summary + "\n\n算法原因仅说明可能机制，不包含算法优化方向。",
    )
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output
